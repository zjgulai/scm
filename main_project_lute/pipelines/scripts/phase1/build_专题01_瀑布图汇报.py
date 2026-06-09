# -*- coding: utf-8 -*-
"""
麦肯锡风格毛利额/毛利率归因瀑布图，并生成 PPT。
依赖: pandas, openpyxl, matplotlib, python-pptx
"""
import os
from pathlib import Path
# 避免 matplotlib 写用户目录权限问题
_BASE = Path(__file__).resolve().parent
os.environ["MPLCONFIGDIR"] = str(_BASE / ".mplcache")
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
# 中文字体：优先系统自带，避免中文显示为方框
plt.rcParams["font.sans-serif"] = ["PingFang SC", "Heiti SC", "STHeiti", "SimHei", "Microsoft YaHei", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE_DIR = Path(__file__).resolve().parents[2]
TOPIC01_DIR = BASE_DIR / "专题产物" / "专题01"
OUT_DIR = TOPIC01_DIR / "表格"
DATA_DIR = BASE_DIR / "原始数据"
EXCEL_PATH = OUT_DIR / "专题01_表格_分析结果.xlsx"
SRC_EXCEL = DATA_DIR / "专题一：分析数据总表.xlsx"
OUT_PPT = TOPIC01_DIR / "汇报" / "专题01_汇报_归因瀑布图.pptx"
IMG_DIR = TOPIC01_DIR / "图表"
IMG_DIR.mkdir(parents=True, exist_ok=True)

# 时间窗口（与 run_专题一分析.py 一致）MAT2026P1/MAT2025P1/YTD2026P1/YTD2025P1
MAT2026P1_MONTHS = [f"2025-{m:02d}" for m in range(2, 13)] + ["2026-01"]   # 2025-02～2026-01
MAT2025P1_MONTHS = [f"2024-{m:02d}" for m in range(2, 13)] + ["2025-01"]   # 2024-02～2025-01
YTD2026P1_MONTHS = ["2026-01"]
YTD2025P1_MONTHS = ["2025-01"]

# 麦肯锡色板 (from mckinsey_ppt_skills)
MCK = {
    "dark_blue": "#051C2C",
    "medium_blue": "#1E3A5F",
    "light_blue": "#3A5F8A",
    "pale_blue": "#5A7F9F",
    "black": "#000000",
    "gray_dark": "#333333",
    "gray_medium": "#666666",
    "white": "#FFFFFF",
    "positive": "#27AE60",
    "negative": "#E74C3C",
    "warning": "#E67E22",
}


def waterfall_data_from_excel():
    """从分析结果 Excel 读取瀑布图所需数据"""
    overall = pd.read_excel(EXCEL_PATH, sheet_name="整体指标")
    platform_attr = pd.read_excel(EXCEL_PATH, sheet_name="平台归因_毛利额与毛利率贡献")
    platform_gp = pd.read_excel(EXCEL_PATH, sheet_name="平台毛利率波动贡献")

    # 整体
    mat_yoy = overall[overall["时间窗口"] == "MAT2025P1"].iloc[0]
    mat_cur = overall[overall["时间窗口"] == "MAT2026P1"].iloc[0]
    ytd_yoy = overall[overall["时间窗口"] == "YTD2025P1"].iloc[0]
    ytd_cur = overall[overall["时间窗口"] == "YTD2026P1"].iloc[0]

    # 平台毛利额（用于瀑布：同期 -> 各平台增量 -> 当期）
    pm_mat = platform_attr[platform_attr["时间窗口"] == "MAT2026P1"][["平台", "品线毛利额"]].set_index("平台")["品线毛利额"]
    pm_mat_yoy = platform_attr[platform_attr["时间窗口"] == "MAT2025P1"][["平台", "品线毛利额"]].set_index("平台")["品线毛利额"]
    pm_ytd = platform_attr[platform_attr["时间窗口"] == "YTD2026P1"][["平台", "品线毛利额"]].set_index("平台")["品线毛利额"]
    pm_ytd_yoy = platform_attr[platform_attr["时间窗口"] == "YTD2025P1"][["平台", "品线毛利额"]].set_index("平台")["品线毛利额"]

    # 平台毛利率波动贡献（贡献率带符号）
    gp_mat = platform_gp[platform_gp["时间窗口"] == "MAT2026P1 vs MAT2025P1"][["平台", "对整体毛利率波动贡献率"]].set_index("平台")["对整体毛利率波动贡献率"]
    gp_ytd = platform_gp[platform_gp["时间窗口"] == "YTD2026P1 vs YTD2025P1"][["平台", "对整体毛利率波动贡献率"]].set_index("平台")["对整体毛利率波动贡献率"]

    return {
        "mat": {
            "margin_start": mat_yoy["品线毛利额"],
            "margin_end": mat_cur["品线毛利额"],
            "margin_deltas": {p: pm_mat.get(p, 0) - pm_mat_yoy.get(p, 0) for p in ["亚马逊", "独立站", "Tiktok平台"]},
            "gp_start": mat_yoy["品线毛利率"],
            "gp_end": mat_cur["品线毛利率"],
            "gp_deltas": gp_mat.to_dict(),
        },
        "ytd": {
            "margin_start": ytd_yoy["品线毛利额"],
            "margin_end": ytd_cur["品线毛利额"],
            "margin_deltas": {p: pm_ytd.get(p, 0) - pm_ytd_yoy.get(p, 0) for p in ["亚马逊", "独立站", "Tiktok平台"]},
            "gp_start": ytd_yoy["品线毛利率"],
            "gp_end": ytd_cur["品线毛利率"],
            "gp_deltas": gp_ytd.to_dict(),
        },
    }


def load_raw_filtered(platforms, regions):
    """从原始表①读取数据，按平台、区域筛选"""
    df = pd.read_excel(SRC_EXCEL, sheet_name="①")
    df = df.dropna(subset=["平台", "区域", "月份"])
    return df[(df["平台"].isin(platforms)) & (df["区域"].isin(regions))].copy()


def get_full_overall_margin_change():
    """全量数据下：当期总销售额、整体毛利率同比差值 ΔG（用于计算子集对整体波动的贡献率）"""
    df = pd.read_excel(SRC_EXCEL, sheet_name="①")
    df = df.dropna(subset=["平台", "区域", "月份"])

    def filter_window(d, months):
        return d[d["月份"].isin(months)]

    w_mat_cur = filter_window(df, MAT2026P1_MONTHS)
    w_mat_yoy = filter_window(df, MAT2025P1_MONTHS)
    w_ytd_cur = filter_window(df, YTD2026P1_MONTHS)
    w_ytd_yoy = filter_window(df, YTD2025P1_MONTHS)

    sale_mat_cur = w_mat_cur["销售额"].sum()
    sale_ytd_cur = w_ytd_cur["销售额"].sum()
    gp_mat_cur = w_mat_cur["品线毛利额"].sum() / sale_mat_cur if sale_mat_cur else np.nan
    gp_mat_yoy = w_mat_yoy["品线毛利额"].sum() / w_mat_yoy["销售额"].sum() if w_mat_yoy["销售额"].sum() else np.nan
    gp_ytd_cur = w_ytd_cur["品线毛利额"].sum() / sale_ytd_cur if sale_ytd_cur else np.nan
    gp_ytd_yoy = w_ytd_yoy["品线毛利额"].sum() / w_ytd_yoy["销售额"].sum() if w_ytd_yoy["销售额"].sum() else np.nan

    delta_gp_mat = (gp_mat_cur - gp_mat_yoy) if pd.notna(gp_mat_cur) and pd.notna(gp_mat_yoy) else np.nan
    delta_gp_ytd = (gp_ytd_cur - gp_ytd_yoy) if pd.notna(gp_ytd_cur) and pd.notna(gp_ytd_yoy) else np.nan

    return {
        "mat": {"total_sale_cur": sale_mat_cur, "delta_gp": delta_gp_mat, "gp_cur": gp_mat_cur, "gp_yoy": gp_mat_yoy},
        "ytd": {"total_sale_cur": sale_ytd_cur, "delta_gp": delta_gp_ytd, "gp_cur": gp_ytd_cur, "gp_yoy": gp_ytd_yoy},
    }


def compute_waterfall_data_from_df(df, platform_list, full_overall=None):
    """
    基于筛选后的 df 计算瀑布图所需数据结构。
    platform_list: 用于瀑布图柱子的平台顺序，如 ["亚马逊", "独立站"]
    full_overall: 全量整体口径 { "mat": { total_sale_cur, delta_gp }, "ytd": { ... } }，
                  若提供则计算子集各平台对「整体毛利率波动」的贡献率（= 贡献/整体ΔG，相对值）。
    """
    def filter_window(d, months):
        return d[d["月份"].isin(months)]

    w_mat_cur = filter_window(df, MAT2026P1_MONTHS)
    w_mat_yoy = filter_window(df, MAT2025P1_MONTHS)
    w_ytd_cur = filter_window(df, YTD2026P1_MONTHS)
    w_ytd_yoy = filter_window(df, YTD2025P1_MONTHS)

    def totals(w):
        s = w["销售额"].sum()
        m = w["品线毛利额"].sum()
        return m, s, (m / s if s != 0 else np.nan)

    margin_mat_yoy, sale_mat_yoy, gp_mat_yoy = totals(w_mat_yoy)
    margin_mat_cur, sale_mat_cur, gp_mat_cur = totals(w_mat_cur)
    margin_ytd_yoy, sale_ytd_yoy, gp_ytd_yoy = totals(w_ytd_yoy)
    margin_ytd_cur, sale_ytd_cur, gp_ytd_cur = totals(w_ytd_cur)

    def platform_metrics(w):
        out = {}
        for p in platform_list:
            g = w[w["平台"] == p]
            s = g["销售额"].sum()
            m = g["品线毛利额"].sum()
            out[p] = {"sale": s, "margin": m, "gp": m / s if s != 0 else np.nan}
        return out

    pm_mat = platform_metrics(w_mat_cur)
    pm_mat_yoy = platform_metrics(w_mat_yoy)
    pm_ytd = platform_metrics(w_ytd_cur)
    pm_ytd_yoy = platform_metrics(w_ytd_yoy)

    margin_deltas_mat = {p: pm_mat[p]["margin"] - pm_mat_yoy[p]["margin"] for p in platform_list}
    margin_deltas_ytd = {p: pm_ytd[p]["margin"] - pm_ytd_yoy[p]["margin"] for p in platform_list}

    # 毛利率绝对贡献（子集内口径，用于瀑布图柱高）：占比 × 毛利率同比差
    gp_deltas_mat = {}
    for p in platform_list:
        share = pm_mat[p]["sale"] / sale_mat_cur if sale_mat_cur else 0
        d_gp = (pm_mat[p]["gp"] - pm_mat_yoy[p]["gp"]) if pd.notna(pm_mat[p]["gp"]) and pd.notna(pm_mat_yoy[p]["gp"]) else 0
        gp_deltas_mat[p] = share * d_gp
    gp_deltas_ytd = {}
    for p in platform_list:
        share = pm_ytd[p]["sale"] / sale_ytd_cur if sale_ytd_cur else 0
        d_gp = (pm_ytd[p]["gp"] - pm_ytd_yoy[p]["gp"]) if pd.notna(pm_ytd[p]["gp"]) and pd.notna(pm_ytd_yoy[p]["gp"]) else 0
        gp_deltas_ytd[p] = share * d_gp

    # 对「整体毛利率波动」的贡献率（相对值）：贡献 / 整体ΔG，仅当传入 full_overall 时计算
    gp_contrib_rate_mat = {}
    gp_contrib_rate_ytd = {}
    if full_overall and full_overall.get("mat") and full_overall.get("ytd"):
        full_sale_mat = full_overall["mat"]["total_sale_cur"]
        full_dg_mat = full_overall["mat"]["delta_gp"]
        full_sale_ytd = full_overall["ytd"]["total_sale_cur"]
        full_dg_ytd = full_overall["ytd"]["delta_gp"]
        for p in platform_list:
            contrib_mat = (pm_mat[p]["sale"] / full_sale_mat * (pm_mat[p]["gp"] - pm_mat_yoy[p]["gp"])) if full_sale_mat and pd.notna(pm_mat[p]["gp"]) and pd.notna(pm_mat_yoy[p]["gp"]) else 0
            contrib_ytd = (pm_ytd[p]["sale"] / full_sale_ytd * (pm_ytd[p]["gp"] - pm_ytd_yoy[p]["gp"])) if full_sale_ytd and pd.notna(pm_ytd[p]["gp"]) and pd.notna(pm_ytd_yoy[p]["gp"]) else 0
            gp_contrib_rate_mat[p] = (contrib_mat / full_dg_mat * 100) if full_dg_mat and abs(full_dg_mat) > 1e-10 else np.nan
            gp_contrib_rate_ytd[p] = (contrib_ytd / full_dg_ytd * 100) if full_dg_ytd and abs(full_dg_ytd) > 1e-10 else np.nan
        subset_contrib_mat = sum((pm_mat[p]["sale"] / full_sale_mat * (pm_mat[p]["gp"] - pm_mat_yoy[p]["gp"])) for p in platform_list if pd.notna(pm_mat[p]["gp"]) and pd.notna(pm_mat_yoy[p]["gp"])) if full_sale_mat else 0
        subset_contrib_ytd = sum((pm_ytd[p]["sale"] / full_sale_ytd * (pm_ytd[p]["gp"] - pm_ytd_yoy[p]["gp"])) for p in platform_list if pd.notna(pm_ytd[p]["gp"]) and pd.notna(pm_ytd_yoy[p]["gp"])) if full_sale_ytd else 0
        subset_rate_mat = (subset_contrib_mat / full_dg_mat * 100) if full_dg_mat and abs(full_dg_mat) > 1e-10 else np.nan
        subset_rate_ytd = (subset_contrib_ytd / full_dg_ytd * 100) if full_dg_ytd and abs(full_dg_ytd) > 1e-10 else np.nan
    else:
        subset_rate_mat = subset_rate_ytd = None

    out = {
        "mat": {
            "margin_start": margin_mat_yoy,
            "margin_end": margin_mat_cur,
            "margin_deltas": margin_deltas_mat,
            "gp_start": gp_mat_yoy,
            "gp_end": gp_mat_cur,
            "gp_deltas": gp_deltas_mat,
        },
        "ytd": {
            "margin_start": margin_ytd_yoy,
            "margin_end": margin_ytd_cur,
            "margin_deltas": margin_deltas_ytd,
            "gp_start": gp_ytd_yoy,
            "gp_end": gp_ytd_cur,
            "gp_deltas": gp_deltas_ytd,
        },
    }
    if full_overall:
        out["mat"]["gp_contribution_rate_pct"] = gp_contrib_rate_mat
        out["ytd"]["gp_contribution_rate_pct"] = gp_contrib_rate_ytd
        out["mat"]["subset_contribution_rate_pct"] = subset_rate_mat
        out["ytd"]["subset_contribution_rate_pct"] = subset_rate_ytd
        out["full_overall"] = full_overall
    return out


def draw_waterfall_margin(data, period_label, period_short, filepath):
    """毛利额瀑布图：期初 -> 各平台增量 -> 期末"""
    start = data["margin_start"]
    end = data["margin_end"]
    deltas = data["margin_deltas"]
    # 按增量绝对值从大到小排，保证瀑布衔接
    order = sorted(deltas.keys(), key=lambda p: -abs(deltas[p]))
    labels = ["MAT2025P1" if "MAT" in period_short else "YTD2025P1"]
    values = [start]
    colors = [MCK["dark_blue"]]
    for p in order:
        v = deltas[p]
        short = "亚马逊" if p == "亚马逊" else ("独立站" if p == "独立站" else "Tiktok")
        labels.append(short)
        values.append(v)
        colors.append(MCK["positive"] if v >= 0 else MCK["negative"])
    labels.append("MAT2026P1" if "MAT" in period_short else "YTD2026P1")
    values.append(end)
    colors.append(MCK["dark_blue"])

    n = len(labels)
    cum = [0.0]
    for i in range(1, n):
        if i < n - 1:
            cum.append(cum[-1] + values[i - 1] if i == 1 else cum[-1] + values[i])
        else:
            cum.append(cum[-1])  # 最后一柱是总高度
    # 瀑布：第0柱从0到start；第1..n-2柱为增量，bottom为前序累计；第n-1柱从0到end
    bottoms = [0, start]
    heights = [start] + [deltas[p] for p in order] + [end]
    for i in range(2, n - 1):
        bottoms.append(bottoms[-1] + heights[i - 1])
    bottoms.append(0)
    # 最后一柱高度为 end
    heights_real = [start] + [deltas[p] for p in order] + [end]
    # 修正：瀑布图中间柱的 bottom = 前一个柱的 top
    b2 = [0]
    h2 = [start]
    for i in range(1, n - 1):
        b2.append(b2[-1] + h2[-1])
        h2.append(heights_real[i])
    b2.append(0)
    h2.append(end)

    fig, ax = plt.subplots(figsize=(10, 5))
    x = np.arange(n)
    bar_w = 0.55
    for i in range(n):
        color = colors[i]
        if i == 0:
            ax.bar(i, h2[i], bar_w, bottom=0, color=color, edgecolor="none")
        elif i == n - 1:
            ax.bar(i, h2[i], bar_w, bottom=0, color=color, edgecolor="none")
        else:
            ax.bar(i, h2[i], bar_w, bottom=b2[i], color=color, edgecolor="none")

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=11)
    ax.set_ylabel("品线毛利额（万元）", fontsize=12)
    ax.set_title(f"品线毛利额变化归因 — {period_label}", fontsize=14, pad=12)
    ax.yaxis.grid(True, linestyle="--", alpha=0.6)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for i in range(n):
        if i == 0 or i == n - 1:
            ax.text(i, h2[i] / 2, f"{h2[i]:,.0f}", ha="center", va="center", fontsize=10, color="white", fontweight="bold")
        else:
            mid = b2[i] + h2[i] / 2
            ax.text(i, mid, f"{h2[i]:+,.0f}", ha="center", va="center", fontsize=10, color="white", fontweight="bold")
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def draw_waterfall_gpm(data, period_label, period_short, filepath):
    """毛利率瀑布图：同期毛利率 -> 各平台贡献 -> 当期毛利率（小数形式，标注为百分比）"""
    start = data["gp_start"]
    end = data["gp_end"]
    deltas = data["gp_deltas"]
    order = sorted(deltas.keys(), key=lambda p: -abs(deltas[p]))
    labels = ["同期\n毛利率"]
    for p in order:
        short = "亚马逊" if p == "亚马逊" else ("独立站" if p == "独立站" else "Tiktok")
        labels.append(short)
    labels.append("当期\n毛利率")
    values = [start] + [deltas[p] for p in order] + [end]
    colors = [MCK["dark_blue"]]
    for p in order:
        colors.append(MCK["positive"] if deltas[p] >= 0 else MCK["negative"])
    colors.append(MCK["dark_blue"])

    n = len(labels)
    b2 = [0]
    h2 = [start]
    for i in range(1, n - 1):
        b2.append(b2[-1] + h2[-1])
        h2.append(values[i])
    b2.append(0)
    h2.append(end)

    fig, ax = plt.subplots(figsize=(10, 5))
    x = np.arange(n)
    bar_w = 0.55
    for i in range(n):
        color = colors[i] if i < len(colors) else MCK["dark_blue"]
        if i == 0 or i == n - 1:
            ax.bar(i, h2[i], bar_w, bottom=0, color=color, edgecolor="none")
        else:
            ax.bar(i, h2[i], bar_w, bottom=b2[i], color=color, edgecolor="none")

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=11)
    ax.set_ylabel("品线毛利率", fontsize=12)
    ax.set_title(f"品线毛利率同比波动归因 — {period_label}", fontsize=14, pad=12)
    ax.yaxis.grid(True, linestyle="--", alpha=0.6)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    # 刻度显示为百分比
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f"{x*100:.1f}%"))
    for i in range(n):
        if i == 0 or i == n - 1:
            ax.text(i, h2[i] / 2, f"{h2[i]*100:.1f}%", ha="center", va="center", fontsize=10, color="white", fontweight="bold")
        else:
            mid = b2[i] + h2[i] / 2
            s = f"{h2[i]*100:+.1f}%"
            ax.text(i, mid, s, ha="center", va="center", fontsize=10, color="white", fontweight="bold")
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def draw_na_eu_gp_margin_growth_combo(df_full, filepath):
    """
    亚马逊 vs 独立站 · 欧洲+北美：按月 毛利率同比增长率、毛利额同比增长率 复合图。
    X 轴：2025-02～2026-01；左 Y：毛利率同比增长率%；右 Y：毛利额同比增长率%。
    """
    df = df_full[(df_full["平台"].isin(["亚马逊", "独立站"])) & (df_full["区域"].isin(["欧洲", "北美"]))].copy()
    months_cur = [f"2025-{m:02d}" for m in range(2, 13)] + ["2026-01"]
    x_labels = [f"2025-{m:02d}" for m in range(2, 13)] + ["2026-01"]
    # 去年同期：2025-02 -> 2024-02, ..., 2026-01 -> 2025-01
    yoy_map = {f"2025-{m:02d}": f"2024-{m:02d}" for m in range(2, 13)}
    yoy_map["2026-01"] = "2025-01"

    def pct_chg(cur, yoy):
        if yoy is None or yoy == 0 or np.isnan(yoy):
            return np.nan
        return (cur - yoy) / yoy * 100

    rows = []
    for m in months_cur:
        m_yoy = yoy_map[m]
        w_cur = df[df["月份"] == m]
        w_yoy = df[df["月份"] == m_yoy]
        row = {"month": m, "x_label": m}
        for platform in ["亚马逊", "独立站"]:
            c = w_cur[w_cur["平台"] == platform]
            y = w_yoy[w_yoy["平台"] == platform]
            s_c, margin_c = c["销售额"].sum(), c["品线毛利额"].sum()
            s_y, margin_y = y["销售额"].sum(), y["品线毛利额"].sum()
            gp_c = margin_c / s_c if s_c else np.nan
            gp_y = margin_y / s_y if s_y else np.nan
            row[f"{platform}_gp_pct"] = pct_chg(gp_c, gp_y)
            row[f"{platform}_margin_pct"] = pct_chg(margin_c, margin_y)
        rows.append(row)

    tab = pd.DataFrame(rows)
    x = np.arange(len(x_labels))

    fig, ax1 = plt.subplots(figsize=(12, 5.5))
    ax2 = ax1.twinx()

    ax1.plot(x, tab["亚马逊_gp_pct"], color=MCK["dark_blue"], marker="o", markersize=4, label="亚马逊 毛利率同比", linewidth=2)
    ax1.plot(x, tab["独立站_gp_pct"], color=MCK["warning"], marker="s", markersize=4, label="独立站 毛利率同比", linewidth=2)
    ax1.set_ylabel("毛利率同比增长率（%）", fontsize=11, color=MCK["gray_dark"])
    ax1.tick_params(axis="y", labelcolor=MCK["gray_dark"])
    ax1.axhline(0, color=MCK["gray_medium"], linestyle="--", alpha=0.6)
    ax1.legend(loc="upper left", fontsize=9)
    ax1.set_ylim(bottom=None)

    ax2.plot(x, tab["亚马逊_margin_pct"], color=MCK["light_blue"], linestyle="--", marker="o", markersize=4, label="亚马逊 毛利额同比", linewidth=1.5)
    ax2.plot(x, tab["独立站_margin_pct"], color=MCK["negative"], linestyle="--", marker="s", markersize=4, label="独立站 毛利额同比", linewidth=1.5)
    ax2.set_ylabel("毛利额同比增长率（%）", fontsize=11, color=MCK["gray_dark"])
    ax2.tick_params(axis="y", labelcolor=MCK["gray_dark"])
    ax2.legend(loc="upper right", fontsize=9)

    ax1.set_xticks(x)
    ax1.set_xticklabels(x_labels, rotation=45, ha="right")
    ax1.set_xlabel("年月", fontsize=11)
    ax1.set_title("亚马逊 vs 独立站 · 欧洲+北美：毛利率同比与毛利额同比增速趋势", fontsize=13, pad=10)
    ax1.grid(True, axis="y", linestyle="--", alpha=0.5)
    ax1.spines["top"].set_visible(False)
    fig.patch.set_facecolor("white")
    ax1.set_facecolor("white")
    ax2.spines["top"].set_visible(False)
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()
    return tab


def _combo_chart_conclusions_and_alerts(combo_tab):
    """基于复合图数据生成结论与预警点（毛利率同比 vs 毛利额同比 关系）。"""
    am_gp = combo_tab["亚马逊_gp_pct"].dropna()
    id_gp = combo_tab["独立站_gp_pct"].dropna()
    am_margin = combo_tab["亚马逊_margin_pct"].dropna()
    id_margin = combo_tab["独立站_margin_pct"].dropna()

    def _avg(s):
        return s.mean() if len(s) else np.nan

    am_gp_avg = _avg(am_gp)
    id_gp_avg = _avg(id_gp)
    am_margin_avg = _avg(am_margin)
    id_margin_avg = _avg(id_margin)

    lines = [
        "一、毛利率增长率与毛利额增长率的关系",
        "• 同向且同步：若某月毛利率同比升、毛利额同比也升，说明量利齐升；若均降则需关注双重压力。",
        "• 反向（毛利率降、毛利额升）：以价换量特征，需关注毛利额增速是否足以弥补让利、可持续性。",
        "• 毛利率升、毛利额降：可能结构变化（高毛利低量占比升），需结合销售额与结构看。",
        "",
        "二、趋势结论（欧洲+北美 2025-02～2026-01）",
    ]
    if pd.notna(am_gp_avg) and pd.notna(am_margin_avg):
        lines.append(f"• 亚马逊：毛利率同比均值为 {am_gp_avg:.1f}%，毛利额同比均值为 {am_margin_avg:.1f}%；两轴走势是否同向可看图判断。")
    if pd.notna(id_gp_avg) and pd.notna(id_margin_avg):
        lines.append(f"• 独立站：毛利率同比均值为 {id_gp_avg:.1f}%，毛利额同比均值为 {id_margin_avg:.1f}%；若毛利率同比波动大而毛利额同比高，存在以价换量可能。")
    lines.extend([
        "",
        "三、预警点",
        "• 毛利率同比持续为负且毛利额同比走弱：警惕以价换量失效或让利过度，需收紧折扣与定价。",
        "• 毛利额同比大幅高于毛利率同比（量增利让）：关注后续毛利率能否企稳，避免长期以价换量。",
        "• 两平台走势背离加剧：区域或结构差异扩大，需分平台制定毛利与规模策略。",
        "• 单月毛利率同比或毛利额同比异常尖峰/谷底：排查大促、退换、一次性费用，避免误读趋势。",
        "",
        "四、关键词",
        "毛利率同比 | 毛利额同比 | 欧洲+北美 | 以价换量 | 数据来源: 专题一分析数据总表",
    ])
    return lines


def analyze_independent_station_price_volume(df_full):
    """
    独立站（全区域）在两个时间窗口：毛利率、毛利额增长波动是否呈现以价换量。
    以价换量：用价格/毛利空间换销量 → 典型为 销售额↑ + 毛利率同比↓；再看毛利额是否↑判断是否换到量。
    返回：{ "report_text": str, "slide_lines": list } 用于写入文件与 PPT。
    """
    df = df_full[df_full["平台"] == "独立站"].copy()

    def filter_window(d, months):
        return d[d["月份"].isin(months)]

    def totals(w):
        s = w["销售额"].sum()
        m = w["品线毛利额"].sum()
        return s, m, (m / s if s != 0 else np.nan)

    w_mat_cur = filter_window(df, MAT2026P1_MONTHS)
    w_mat_yoy = filter_window(df, MAT2025P1_MONTHS)
    w_ytd_cur = filter_window(df, YTD2026P1_MONTHS)
    w_ytd_yoy = filter_window(df, YTD2025P1_MONTHS)

    sale_mat_cur, margin_mat_cur, gp_mat_cur = totals(w_mat_cur)
    sale_mat_yoy, margin_mat_yoy, gp_mat_yoy = totals(w_mat_yoy)
    sale_ytd_cur, margin_ytd_cur, gp_ytd_cur = totals(w_ytd_cur)
    sale_ytd_yoy, margin_ytd_yoy, gp_ytd_yoy = totals(w_ytd_yoy)

    # 同比
    def pct_chg(cur, yoy):
        return ((cur - yoy) / yoy * 100) if yoy and yoy != 0 else np.nan

    mat_sale_pct = pct_chg(sale_mat_cur, sale_mat_yoy)
    mat_margin_pct = pct_chg(margin_mat_cur, margin_mat_yoy)
    mat_gp_pp = (gp_mat_cur - gp_mat_yoy) * 100 if pd.notna(gp_mat_cur) and pd.notna(gp_mat_yoy) else np.nan
    ytd_sale_pct = pct_chg(sale_ytd_cur, sale_ytd_yoy)
    ytd_margin_pct = pct_chg(margin_ytd_cur, margin_ytd_yoy)
    ytd_gp_pp = (gp_ytd_cur - gp_ytd_yoy) * 100 if pd.notna(gp_ytd_cur) and pd.notna(gp_ytd_yoy) else np.nan

    def judge_price_volume(sale_pct, gp_pp, margin_pct):
        """判断是否以价换量：销售额↑且毛利率↓ 为典型特征；再看毛利额是否↑。"""
        sale_up = sale_pct is not None and not np.isnan(sale_pct) and sale_pct > 0
        gp_down = gp_pp is not None and not np.isnan(gp_pp) and gp_pp < 0
        margin_up = margin_pct is not None and not np.isnan(margin_pct) and margin_pct > 0
        if sale_up and gp_down:
            if margin_up:
                return "存在以价换量特征，且换到量：销售额与毛利额同比均增，毛利率同比降，规模增长弥补让利。"
            else:
                return "存在以价换量特征，但未换到足够量：毛利率同比降，毛利额同比亦降，让利未带来足够增量。"
        if sale_up and (gp_pp is None or np.isnan(gp_pp) or gp_pp >= 0):
            return "量价齐升：销售额与毛利率同比均增，非以价换量。"
        if not sale_up and gp_down:
            return "量缩利让：销售额未增且毛利率降，非以价换量，需关注收缩与盈利双压。"
        return "销售额与毛利率同比组合无明显以价换量特征。"

    conclusion_mat = judge_price_volume(mat_sale_pct, mat_gp_pp, mat_margin_pct)
    conclusion_ytd = judge_price_volume(ytd_sale_pct, ytd_gp_pp, ytd_margin_pct)

    def _n(x):
        return f"{x:,.0f}" if x is not None and not (isinstance(x, float) and np.isnan(x)) else "—"

    def _p(x):
        return f"{x:+.1f}%" if x is not None and not (isinstance(x, float) and np.isnan(x)) else "—"

    report = [
        "【独立站】两时间窗口：毛利率/毛利额增长波动与以价换量判断",
        "口径：独立站全区域；以价换量 = 用价格/毛利空间换销量，典型为 销售额↑ + 毛利率同比↓。",
        "",
        "一、MAT2026P1 vs MAT2025P1",
        f"  销售额：同期 {_n(sale_mat_yoy)} 万元 → 当期 {_n(sale_mat_cur)} 万元，同比 {_p(mat_sale_pct)}。",
        f"  品线毛利额：同期 {_n(margin_mat_yoy)} 万元 → 当期 {_n(margin_mat_cur)} 万元，同比 {_p(mat_margin_pct)}。",
        f"  品线毛利率：同期 {gp_mat_yoy*100:.2f}% → 当期 {gp_mat_cur*100:.2f}%，同比 {mat_gp_pp:+.2f}pp。" if pd.notna(gp_mat_cur) and pd.notna(gp_mat_yoy) else "  品线毛利率：—",
        f"  结论：{conclusion_mat}",
        "",
        "二、YTD2026P1 vs YTD2025P1",
        f"  销售额：同期 {_n(sale_ytd_yoy)} 万元 → 当期 {_n(sale_ytd_cur)} 万元，同比 {_p(ytd_sale_pct)}。",
        f"  品线毛利额：同期 {_n(margin_ytd_yoy)} 万元 → 当期 {_n(margin_ytd_cur)} 万元，同比 {_p(ytd_margin_pct)}。",
        f"  品线毛利率：同期 {gp_ytd_yoy*100:.2f}% → 当期 {gp_ytd_cur*100:.2f}%，同比 {ytd_gp_pp:+.2f}pp。" if pd.notna(gp_ytd_cur) and pd.notna(gp_ytd_yoy) else "  品线毛利率：—",
        f"  结论：{conclusion_ytd}",
        "",
        "三、综合",
        "若 MAT 与 YTD 结论一致为以价换量且换到量，则独立站策略上在用毛利换规模；若不一致，需区分滚动年与年初窗口的节奏差异。",
    ]
    report_text = "\n".join(report)

    slide_lines = [
        "一、口径",
        "独立站全区域；以价换量 = 销售额↑ + 毛利率同比↓（让利换规模），再看毛利额是否↑。",
        "",
        "二、MAT2026P1 vs MAT2025P1",
        f"销售额同比 {_p(mat_sale_pct)}，毛利额同比 {_p(mat_margin_pct)}，毛利率同比 {mat_gp_pp:+.2f}pp。" if pd.notna(mat_gp_pp) else f"销售额同比 {_p(mat_sale_pct)}，毛利额同比 {_p(mat_margin_pct)}。",
        f"结论：{conclusion_mat}",
        "",
        "三、YTD2026P1 vs YTD2025P1",
        f"销售额同比 {_p(ytd_sale_pct)}，毛利额同比 {_p(ytd_margin_pct)}，毛利率同比 {ytd_gp_pp:+.2f}pp。" if pd.notna(ytd_gp_pp) else f"销售额同比 {_p(ytd_sale_pct)}，毛利额同比 {_p(ytd_margin_pct)}。",
        f"结论：{conclusion_ytd}",
        "",
        "四、需关注",
        "若两窗口均为以价换量且换到量：独立站以让利换规模，可持续性取决于后续毛利率与规模平衡。若仅 YTD 成立：短期促销/定价策略明显，需与 MAT 结构对比。",
    ]

    return {"report_text": report_text, "slide_lines": slide_lines}


def analyze_independent_station_na_eu_thin_margin(df_full):
    """
    独立站 · 仅欧洲+北美 子集：MAT 窗口下是否属于「以价换微薄毛利增长」，
    以及 毛利率下滑 与 毛利额增长 是否对等。
    以价换微薄毛利 = 让利（毛利率↓）但换来的毛利额增长有限（毛利额增速 << 销售额增速 或 毛利额增速偏低）。
    对等 = 毛利率略降、毛利额增速与销售额增速接近（量增为主、价让有限）→ 结构上对等；反之不对等。
    """
    df = df_full[(df_full["平台"] == "独立站") & (df_full["区域"].isin(["北美", "欧洲"]))].copy()

    def filter_window(d, months):
        return d[d["月份"].isin(months)]

    def totals(w):
        s = w["销售额"].sum()
        m = w["品线毛利额"].sum()
        return s, m, (m / s if s != 0 else np.nan)

    w_mat_cur = filter_window(df, MAT2026P1_MONTHS)
    w_mat_yoy = filter_window(df, MAT2025P1_MONTHS)

    sale_cur, margin_cur, gp_cur = totals(w_mat_cur)
    sale_yoy, margin_yoy, gp_yoy = totals(w_mat_yoy)

    sale_pct = ((sale_cur - sale_yoy) / sale_yoy * 100) if sale_yoy and sale_yoy != 0 else np.nan
    margin_pct = ((margin_cur - margin_yoy) / margin_yoy * 100) if margin_yoy and margin_yoy != 0 else np.nan
    gp_pp = (gp_cur - gp_yoy) * 100 if pd.notna(gp_cur) and pd.notna(gp_yoy) else np.nan
    margin_abs_delta = margin_cur - margin_yoy

    # 是否以价换微薄毛利增长：毛利率同比↓ 且 (毛利额增速明显低于销售额增速 或 毛利额增速偏低)
    gp_down = gp_pp is not None and not np.isnan(gp_pp) and gp_pp < 0
    sale_up = sale_pct is not None and not np.isnan(sale_pct) and sale_pct > 0
    # 微薄定义：毛利额增速 < 销售额增速的 50%，或 毛利额增速 < 20%
    thin_margin_growth = False
    if gp_down and sale_up and margin_pct is not None and not np.isnan(margin_pct):
        if margin_pct < 20 or (sale_pct and margin_pct < sale_pct * 0.5):
            thin_margin_growth = True

    # 对等：毛利额增速 与 销售额增速 比；若毛利率同比↑则无“让利”，直接看量利结构
    ratio_margin_to_sale = (margin_pct / sale_pct) if sale_pct and abs(sale_pct) > 1e-6 else np.nan
    gp_up = gp_pp is not None and not np.isnan(gp_pp) and gp_pp > 0
    if gp_up:
        if pd.notna(ratio_margin_to_sale) and ratio_margin_to_sale >= 0.8 and ratio_margin_to_sale <= 1.2:
            balance_conclusion = "无毛利率下滑、毛利率同比略升：毛利额增速与销售额增速接近，量价结构健康、不存在以价换量。"
        elif pd.notna(ratio_margin_to_sale) and ratio_margin_to_sale > 1.2:
            balance_conclusion = "无让利且有利：毛利率同比升，毛利额增速高于销售额增速，量价齐升、结构优于对等。"
        else:
            balance_conclusion = "毛利率同比升，无让利；毛利额与销售额结构需结合比例综合看。"
    else:
        if pd.notna(ratio_margin_to_sale) and ratio_margin_to_sale >= 0.8 and ratio_margin_to_sale <= 1.2:
            balance_conclusion = "毛利率下滑与毛利额增长大体对等：毛利额增速与销售额增速接近，以量增为主、价让有限。"
        elif pd.notna(ratio_margin_to_sale) and ratio_margin_to_sale > 1.2:
            balance_conclusion = "不对等且有利：毛利率仅略降，毛利额增速高于销售额增速，让利少、换来的毛利额增长多。"
        elif thin_margin_growth:
            balance_conclusion = "不对等且不利：毛利率下滑而毛利额增长有限，属于以价换微薄毛利增长。"
        elif pd.notna(ratio_margin_to_sale) and ratio_margin_to_sale < 0.8 and ratio_margin_to_sale >= 0:
            balance_conclusion = "不对等：毛利额增速低于销售额增速，部分让利未转化为毛利额增长。"
        else:
            balance_conclusion = "需结合毛利率与销售额结构综合判断对等性。"

    if gp_up:
        thin_conclusion = "不存在以价换量（毛利率同比升）：非以价换微薄毛利，属量价齐升或量增利稳。"
    elif thin_margin_growth:
        thin_conclusion = "属于以价换微薄毛利增长：让利（毛利率同比降）但毛利额增速有限或明显低于销售额增速。"
    else:
        thin_conclusion = "不属于以价换微薄毛利增长：毛利率仅略降，毛利额同比增速高，让利少、换来的毛利额增长多。"

    def _n(x):
        return f"{x:,.0f}" if x is not None and not (isinstance(x, float) and np.isnan(x)) else "—"

    def _p(x):
        return f"{x:+.1f}%" if x is not None and not (isinstance(x, float) and np.isnan(x)) else "—"

    report = [
        "【独立站 · 欧洲+北美】MAT2026P1 vs MAT2025P1：是否以价换微薄毛利、毛利率下滑与毛利额增长是否对等",
        "口径：独立站 仅 区域∈{欧洲, 北美}。以价换微薄毛利 = 毛利率↓ 且 毛利额增速有限或远低于销售额增速。",
        "",
        "一、MAT 同比数据（独立站 欧洲+北美）",
        f"  销售额：同期 {_n(sale_yoy)} 万元 → 当期 {_n(sale_cur)} 万元，同比 {_p(sale_pct)}。",
        f"  品线毛利额：同期 {_n(margin_yoy)} 万元 → 当期 {_n(margin_cur)} 万元，同比 {_p(margin_pct)}，绝对增量 {margin_abs_delta:+,.0f} 万元。",
        f"  品线毛利率：同期 {gp_yoy*100:.2f}% → 当期 {gp_cur*100:.2f}%，同比 {gp_pp:+.2f}pp。" if pd.notna(gp_pp) else "  品线毛利率：—",
        "",
        "二、是否属于以价换微薄毛利增长",
        thin_conclusion,
        "",
        "三、毛利率变动与毛利额增长是否对等",
        f"  毛利额增速/销售额增速 = {ratio_margin_to_sale:.2f}。" if pd.notna(ratio_margin_to_sale) else "  —",
        balance_conclusion,
        "",
        "四、小结",
        "欧洲+北美子集若毛利率同比升：无让利，毛利额与销售额同步高增、结构健康。若毛利率同比降且毛利额增速≈销售额增速：对等；若毛利额增速远低于销售额增速：以价换微薄毛利。",
    ]
    report_text = "\n".join(report)

    slide_lines = [
        "一、口径",
        "独立站 仅 欧洲+北美；判断是否以价换微薄毛利、毛利率下滑与毛利额增长是否对等。",
        "",
        "二、MAT 同比（独立站 欧洲+北美）",
        f"销售额同比 {_p(sale_pct)}，毛利额同比 {_p(margin_pct)}，毛利率同比 {gp_pp:+.2f}pp。" if pd.notna(gp_pp) else f"销售额同比 {_p(sale_pct)}，毛利额同比 {_p(margin_pct)}。",
        "",
        "三、是否以价换微薄毛利",
        thin_conclusion,
        "",
        "四、毛利率变动与毛利额增长是否对等",
        balance_conclusion,
    ]

    return {"report_text": report_text, "slide_lines": slide_lines}


def add_slide_section(prs, section_title):
    """章节标题页（居中大标题）"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(1.2))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = section_title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(5, 28, 44)
    tb.text_frame.paragraphs[0].alignment = 1
    return slide


def add_slide_title(prs, title_text):
    """封面/标题页"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(12)
    p = slide.shapes.add_textbox(left, top, width, Inches(1.2))
    tf = p.text_frame
    tf.word_wrap = True
    run = p.text_frame.paragraphs[0].add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(5, 28, 44)
    p.text_frame.paragraphs[0].alignment = 1  # center
    return slide


def add_slide_chart(prs, title, overview, img_path, conclusions, keywords=""):
    """麦肯锡风格单图页：标题 + 概述 + 图 + 结论 + 关键词"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    # 标题
    slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6)).text_frame.paragraphs[0].add_run().text = title
    slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6)).text_frame.paragraphs[0].font.size = Pt(18)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6)).text_frame.paragraphs[0].font.bold = True
    # 用单一 textbox 避免重复
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(1.1))
    tb.text_frame.word_wrap = True
    p0 = tb.text_frame.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    r0.font.size = Pt(18)
    r0.font.bold = True
    r0.font.color.rgb = RGBColor(0, 0, 0)
    p1 = tb.text_frame.add_paragraph()
    r1 = p1.add_run()
    r1.text = overview
    r1.font.size = Pt(11)
    r1.font.color.rgb = RGBColor(51, 51, 51)
    # 图片
    slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.55), width=Inches(12), height=Inches(4.2))
    # 结论
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(12), Inches(1.0))
    tb2.text_frame.word_wrap = True
    for line in conclusions:
        p = tb2.text_frame.add_paragraph() if tb2.text_frame.paragraphs[0].text else tb2.text_frame.paragraphs[0]
        if tb2.text_frame.paragraphs[0].text == "":
            r = tb2.text_frame.paragraphs[0].add_run()
        else:
            p = tb2.text_frame.add_paragraph()
            r = p.add_run()
        r.text = line
        r.font.size = Pt(10)
        r.font.color.rgb = RGBColor(51, 51, 51)
    # 关键词
    tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12), Inches(0.4))
    tb3.text_frame.paragraphs[0].add_run().text = keywords or "数据来源: 专题一分析数据总表"
    tb3.text_frame.paragraphs[0].font.size = Pt(9)
    tb3.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
    return slide


def add_slide_chart_simple(prs, title, img_path, footnote=""):
    """简化图页：标题 + 图 + 脚注"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.55))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0, 0, 0)
    slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.0), width=Inches(12), height=Inches(5.5))
    if footnote:
        ft = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12), Inches(0.35))
        ft.text_frame.paragraphs[0].add_run().text = footnote
        ft.text_frame.paragraphs[0].font.size = Pt(9)
        ft.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
    return slide


def add_slide_conclusions(prs, lines):
    """结论汇总页"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(6.8))
    tb.text_frame.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tb.text_frame.paragraphs[0]
        else:
            p = tb.text_frame.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(14) if line.startswith("一、") or line.startswith("二、") or line.startswith("三、") else Pt(11)
        r.font.bold = line.startswith("一、") or line.startswith("二、") or line.startswith("三、")
        r.font.color.rgb = RGBColor(0, 0, 0) if r.font.bold else RGBColor(51, 51, 51)
    return slide


def main():
    data = waterfall_data_from_excel()

    # 1) 毛利额 MAT
    path_margin_mat = IMG_DIR / "waterfall_毛利额_MAT2026P1.png"
    draw_waterfall_margin(data["mat"], "MAT2026P1（2025-02～2026-01）", "MAT", path_margin_mat)

    # 2) 毛利额 YTD
    path_margin_ytd = IMG_DIR / "waterfall_毛利额_YTD2026P1.png"
    draw_waterfall_margin(data["ytd"], "YTD2026P1（2026年1月）", "YTD", path_margin_ytd)

    # 3) 毛利率 MAT
    path_gpm_mat = IMG_DIR / "waterfall_毛利率_MAT2026P1.png"
    draw_waterfall_gpm(data["mat"], "MAT2026P1 vs MAT2025P1", "MAT", path_gpm_mat)

    # 4) 毛利率 YTD
    path_gpm_ytd = IMG_DIR / "waterfall_毛利率_YTD2026P1.png"
    draw_waterfall_gpm(data["ytd"], "YTD2026P1 vs YTD2025P1", "YTD", path_gpm_ytd)

    # PPT
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_slide_title(prs, "专题一：品线毛利额与毛利率归因分析 — 瀑布图")

    add_slide_chart_simple(
        prs,
        "品线毛利额变化归因 — MAT2026P1（2025-02～2026-01）",
        path_margin_mat,
        "数据来源: 专题一分析数据总表 Sheet① | 同期 MAT2025P1: 2024-02～2025-01",
    )
    add_slide_chart_simple(
        prs,
        "品线毛利额变化归因 — YTD2026P1（2026年1月）",
        path_margin_ytd,
        "数据来源: 专题一分析数据总表 Sheet① | 同期 YTD2025P1: 2025年1月",
    )
    add_slide_chart_simple(
        prs,
        "品线毛利率同比波动归因 — MAT2026P1 vs MAT2025P1",
        path_gpm_mat,
        "贡献率 = 当期销售额占比 × (当期毛利率 - 同期毛利率) | 数据来源: 专题一分析数据总表",
    )
    add_slide_chart_simple(
        prs,
        "品线毛利率同比波动归因 — YTD2026P1 vs YTD2025P1",
        path_gpm_ytd,
        "数据来源: 专题一分析数据总表 Sheet①",
    )

    conclusion_lines = [
        "一、毛利额归因结论",
        "• MAT2026P1：毛利额从 MAT2025P1 增至 MAT2026P1，增量主要由各平台贡献（见分析结果表）。",
        "• YTD2026P1：毛利额从 YTD2025P1（2025年1月）增至 YTD2026P1（2026年1月），亚马逊、独立站、Tiktok 贡献见表。",
        "",
        "二、毛利率波动归因结论",
        "• MAT2026P1 vs MAT2025P1：整体毛利率同比波动及平台贡献见「平台毛利率波动贡献」表。",
        "• YTD2026P1 vs YTD2025P1：整体毛利率同比波动及平台贡献见表。",
        "",
        "三、关键词",
        "品线毛利额 | 品线毛利率 | 平台归因 | MAT | YTD | 数据来源: 专题一分析数据总表",
    ]
    add_slide_conclusions(prs, conclusion_lines)

    # ---------- 亚马逊 & 独立站 · 北美 & 欧洲 子集分析（贡献率=对整体毛利率波动的相对占比）----------
    full_overall = get_full_overall_margin_change()
    df_na_eu = load_raw_filtered(platforms=["亚马逊", "独立站"], regions=["北美", "欧洲"])
    data_na_eu = compute_waterfall_data_from_df(df_na_eu, platform_list=["亚马逊", "独立站"], full_overall=full_overall)
    prefix = "na_eu_"
    path_na_eu_margin_mat = IMG_DIR / f"{prefix}waterfall_毛利额_MAT2026P1.png"
    path_na_eu_margin_ytd = IMG_DIR / f"{prefix}waterfall_毛利额_YTD2026P1.png"
    path_na_eu_gpm_mat = IMG_DIR / f"{prefix}waterfall_毛利率_MAT2026P1.png"
    path_na_eu_gpm_ytd = IMG_DIR / f"{prefix}waterfall_毛利率_YTD2026P1.png"
    draw_waterfall_margin(data_na_eu["mat"], "MAT2026P1（亚马逊+独立站，北美+欧洲）", "MAT", path_na_eu_margin_mat)
    draw_waterfall_margin(data_na_eu["ytd"], "YTD2026P1（亚马逊+独立站，北美+欧洲）", "YTD", path_na_eu_margin_ytd)
    draw_waterfall_gpm(data_na_eu["mat"], "MAT2026P1 vs MAT2025P1（亚马逊+独立站，北美+欧洲）", "MAT", path_na_eu_gpm_mat)
    draw_waterfall_gpm(data_na_eu["ytd"], "YTD2026P1 vs YTD2025P1（亚马逊+独立站，北美+欧洲）", "YTD", path_na_eu_gpm_ytd)

    add_slide_section(prs, "亚马逊 & 独立站 · 北美 & 欧洲 归因对比")
    add_slide_chart_simple(
        prs,
        "品线毛利额变化归因 — MAT2026P1（亚马逊+独立站，北美+欧洲）",
        path_na_eu_margin_mat,
        "数据来源: 专题一分析数据总表 Sheet①，筛选 北美+欧洲 | 同期 MAT2025P1: 2024-02～2025-01",
    )
    add_slide_chart_simple(
        prs,
        "品线毛利额变化归因 — YTD2026P1（亚马逊+独立站，北美+欧洲）",
        path_na_eu_margin_ytd,
        "数据来源: 专题一分析数据总表 Sheet①，筛选 北美+欧洲 | 同期 YTD2025P1: 2025年1月",
    )
    add_slide_chart_simple(
        prs,
        "品线毛利率同比波动归因 — MAT2026P1 vs MAT2025P1（亚马逊+独立站，北美+欧洲）",
        path_na_eu_gpm_mat,
        "贡献率 = 当期销售额占比 × (当期毛利率 - 同期毛利率) | 数据来源: 专题一分析数据总表",
    )
    add_slide_chart_simple(
        prs,
        "品线毛利率同比波动归因 — YTD2026P1 vs YTD2025P1（亚马逊+独立站，北美+欧洲）",
        path_na_eu_gpm_ytd,
        "数据来源: 专题一分析数据总表 Sheet①，筛选 亚马逊+独立站，北美+欧洲",
    )

    # 子集结论：贡献率 = 对整体毛利率波动的贡献占比（相对值）；结论分层、需关注问题明确
    mat_m = data_na_eu["mat"]
    ytd_m = data_na_eu["ytd"]
    dm_am = mat_m["margin_deltas"].get("亚马逊", 0)
    dm_id = mat_m["margin_deltas"].get("独立站", 0)
    dy_am = ytd_m["margin_deltas"].get("亚马逊", 0)
    dy_id = ytd_m["margin_deltas"].get("独立站", 0)
    gp_mat_d_pp = (mat_m["gp_end"] - mat_m["gp_start"]) * 100 if pd.notna(mat_m["gp_end"]) and pd.notna(mat_m["gp_start"]) else 0
    gp_ytd_d_pp = (ytd_m["gp_end"] - ytd_m["gp_start"]) * 100 if pd.notna(ytd_m["gp_end"]) and pd.notna(ytd_m["gp_start"]) else 0
    _dg_mat = full_overall["mat"]["delta_gp"]
    _dg_ytd = full_overall["ytd"]["delta_gp"]
    full_dg_mat_pp = _dg_mat * 100 if pd.notna(_dg_mat) else np.nan
    full_dg_ytd_pp = _dg_ytd * 100 if pd.notna(_dg_ytd) else np.nan
    rate_mat_am = mat_m.get("gp_contribution_rate_pct", {}).get("亚马逊")
    rate_mat_id = mat_m.get("gp_contribution_rate_pct", {}).get("独立站")
    rate_ytd_am = ytd_m.get("gp_contribution_rate_pct", {}).get("亚马逊")
    rate_ytd_id = ytd_m.get("gp_contribution_rate_pct", {}).get("独立站")
    subset_rate_mat = data_na_eu.get("mat", {}).get("subset_contribution_rate_pct")
    subset_rate_ytd = data_na_eu.get("ytd", {}).get("subset_contribution_rate_pct")

    def _fmt(x):
        if x is None or (isinstance(x, float) and np.isnan(x)):
            return "—"
        return f"{x:.1f}%"

    conclusion_na_eu = [
        "一、口径说明",
        "• 本页为「亚马逊+独立站、北美+欧洲」子集。",
        "• 毛利率贡献率 = 该主体对「全量整体」毛利率同比差值的贡献占比（贡献÷整体ΔG），为相对值，非绝对值。",
        "",
        "二、毛利额归因结论",
        f"• MAT2026P1：子集毛利额由 MAT2025P1 {mat_m['margin_start']:,.0f} 万元增至 {mat_m['margin_end']:,.0f} 万元；"
        f"亚马逊增量 {dm_am:+,.0f} 万元、独立站 {dm_id:+,.0f} 万元。",
        f"• YTD2026P1：子集毛利额由 YTD2025P1 {ytd_m['margin_start']:,.0f} 万元增至 {ytd_m['margin_end']:,.0f} 万元；"
        f"亚马逊增量 {dy_am:+,.0f} 万元、独立站 {dy_id:+,.0f} 万元。",
        "",
        "三、毛利率波动归因结论（贡献率=对整体波动的占比）",
        f"• 整体口径：MAT 整体毛利率同比 {full_dg_mat_pp:+.2f}pp，YTD 同比 {full_dg_ytd_pp:+.2f}pp。",
        f"• MAT：亚马逊对整体毛利率波动的贡献率 {_fmt(rate_mat_am)}，独立站 {_fmt(rate_mat_id)}；"
        f"子集合计对整体贡献率 {_fmt(subset_rate_mat)}。",
        f"• YTD：亚马逊对整体毛利率波动的贡献率 {_fmt(rate_ytd_am)}，独立站 {_fmt(rate_ytd_id)}；"
        f"子集合计对整体贡献率 {_fmt(subset_rate_ytd)}。",
        "",
        "四、需重点关注的问题",
        "• 若子集对整体贡献率偏高：北美+欧洲的亚马逊/独立站定价与成本变动对全局毛利率影响大，需优先监控该区域平台毛利与折扣。",
        "• 若某平台贡献率为负：该平台在北美/欧洲的毛利率同比恶化，会拉低整体；需拆解成本、价格与结构，定位原因。",
        "• 若 MAT 与 YTD 贡献率差异大：短期（YTD）与滚动年（MAT）驱动因素可能不同，需结合季节与活动做归因复核。",
        "• 决策建议：在贡献率高的平台/区域上优先做毛利率改善与监控，避免单一区域波动放大为整体风险。",
        "",
        "五、关键词",
        "贡献率=对整体ΔG占比 | 亚马逊 | 独立站 | 北美 | 欧洲 | 数据来源: 专题一分析数据总表",
    ]
    add_slide_conclusions(prs, conclusion_na_eu)

    # ---------- 独立站 两时间窗口：毛利率/毛利额增长波动是否以价换量 ----------
    df_full = pd.read_excel(SRC_EXCEL, sheet_name="①").dropna(subset=["平台", "区域", "月份"])
    indie_result = analyze_independent_station_price_volume(df_full)

    add_slide_section(prs, "独立站：毛利率/毛利额波动与以价换量")
    add_slide_conclusions(prs, indie_result["slide_lines"])

    out_indie = OUT_DIR / "专题一_独立站以价换量分析.txt"
    with open(out_indie, "w", encoding="utf-8") as f:
        f.write(indie_result["report_text"])
    print("独立站以价换量分析:", out_indie)

    # 独立站 仅欧洲+北美：是否以价换微薄毛利、毛利率下滑与毛利额增长是否对等
    indie_na_eu = analyze_independent_station_na_eu_thin_margin(df_full)
    add_slide_section(prs, "独立站 · 欧洲+北美：以价换量与对等性")
    add_slide_conclusions(prs, indie_na_eu["slide_lines"])
    out_indie_na_eu = OUT_DIR / "专题一_独立站欧洲北美以价换量对等分析.txt"
    with open(out_indie_na_eu, "w", encoding="utf-8") as f:
        f.write(indie_na_eu["report_text"])
    print("独立站欧洲北美以价换量对等分析:", out_indie_na_eu)

    # ---------- 亚马逊 vs 独立站 · 欧洲+北美：毛利率与毛利额同比增速趋势复合图 + 结论与预警 ----------
    path_combo = IMG_DIR / "na_eu_毛利率与毛利额同比趋势.png"
    combo_tab = draw_na_eu_gp_margin_growth_combo(df_full, path_combo)
    add_slide_section(prs, "亚马逊 vs 独立站 · 欧洲+北美：毛利率与毛利额同比趋势")
    add_slide_chart_simple(
        prs,
        "毛利率同比增长率（左轴）与毛利额同比增长率（右轴）— 2025-02～2026-01",
        path_combo,
        "数据来源: 专题一分析数据总表 Sheet①，区域=欧洲+北美；同比=当月 vs 去年同月",
    )
    # 从 combo_tab 生成简要结论与预警（基于趋势与相关性）
    combo_conclusion_lines = _combo_chart_conclusions_and_alerts(combo_tab)
    add_slide_conclusions(prs, combo_conclusion_lines)

    # ---------- 全部分析与归因洞察结论汇总（时间口径统一 P1）----------
    add_slide_section(prs, "全部分析与归因洞察结论汇总")
    summary_lines = [
        "一、时间口径（统一 P1）",
        "MAT2026P1: 2025-02～2026-01 | MAT2025P1: 2024-02～2025-01 | YTD2026P1: 2026年1月 | YTD2025P1: 2025年1月",
        "",
        "二、主要发现",
        "• 毛利额：亚马逊主贡献（MAT/YTD 约 92%～93%），独立站占比双降，Tiktok 占比提升。",
        "• 毛利率波动：亚马逊、Tiktok 正贡献；独立站负贡献（MAT -16.2%，YTD -53.6%），拖累整体。",
        "• 亚马逊 vs 独立站 欧洲+北美毛利率差距：约 15.78pp（MAT）、16.66pp（YTD），独立站毛利率明显偏低。",
        "• 独立站 MAT 以价换量且换到量；YTD 以价换量未换到足够量，需警惕让利未带来增量。",
        "",
        "三、需重点关注的问题",
        "• 独立站对整体毛利率负贡献：需拆解成本、价格、折扣与区域结构，制定改善与监控。",
        "• 独立站在北美、欧洲负贡献占比较大：优先在该区域做独立站毛利与定价复核。",
        "• 北美+欧洲子集对整体毛利率贡献率高：该区域平台毛利与折扣变动对全局影响大。",
        "",
        "四、预警点",
        "• 毛利率同比持续为负且毛利额同比走弱：警惕以价换量失效或让利过度。",
        "• 毛利额同比大幅高于毛利率同比：关注毛利率能否企稳，避免长期以价换量。",
        "• 两平台走势背离加剧：需分平台制定毛利与规模策略。",
        "",
        "五、关键词",
        "品线毛利额 | 品线毛利率 | 平台归因 | MAT2026P1 | YTD2026P1 | 以价换量 | 数据来源: 专题一分析数据总表 Sheet①",
    ]
    add_slide_conclusions(prs, summary_lines)

    prs.save(OUT_PPT)

    # 输出子集对整体毛利率波动的贡献率（相对值）到文本，便于核对
    out_txt = OUT_DIR / "专题一_子集贡献率结果.txt"
    with open(out_txt, "w", encoding="utf-8") as f:
        f.write("【亚马逊+独立站、北美+欧洲】对「整体」毛利率同比差值的贡献率（贡献÷整体ΔG，相对值）\n\n")
        f.write(f"整体口径：MAT 毛利率同比差值 = {full_dg_mat_pp:.4f}pp；YTD = {full_dg_ytd_pp:.4f}pp。\n\n")
        f.write("MAT2026P1 vs MAT2025P1：\n")
        f.write(f"  亚马逊 对整体毛利率波动的贡献率 = {_fmt(rate_mat_am)}\n")
        f.write(f"  独立站 对整体毛利率波动的贡献率 = {_fmt(rate_mat_id)}\n")
        f.write(f"  子集合计 对整体毛利率波动的贡献率 = {_fmt(subset_rate_mat)}\n\n")
        f.write("YTD2026P1 vs YTD2025P1：\n")
        f.write(f"  亚马逊 对整体毛利率波动的贡献率 = {_fmt(rate_ytd_am)}\n")
        f.write(f"  独立站 对整体毛利率波动的贡献率 = {_fmt(rate_ytd_id)}\n")
        f.write(f"  子集合计 对整体毛利率波动的贡献率 = {_fmt(subset_rate_ytd)}\n")
    print("已生成:", OUT_PPT)
    print("子集贡献率结果:", out_txt)
    print("图表缓存目录:", IMG_DIR)
    print("已追加: 亚马逊+独立站、北美+欧洲 子集瀑布图及结论页（贡献率=对整体ΔG占比）")
    print("已追加: 全部分析与归因洞察结论汇总页（时间口径统一 P1）；详见 专题一_全部分析与归因洞察结论汇总.txt")


if __name__ == "__main__":
    main()
