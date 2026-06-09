# MOVED TO: scripts/phase1/build_专题04_四因素瀑布图.py
# -*- coding: utf-8 -*-
"""
专题四 Sheet④：客品数、品单价、订单数、前台毛利率 对整体毛利额同比波动的贡献占比瀑布图，追加到 专题一_归因瀑布图.pptx
"""
import os
from pathlib import Path
_BASE = Path(__file__).resolve().parent
os.environ["MPLCONFIGDIR"] = str(_BASE / ".mplcache")
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
plt.rcParams["font.sans-serif"] = ["PingFang SC", "Heiti SC", "STHeiti", "SimHei", "Microsoft YaHei", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE_DIR = Path(__file__).resolve().parents[2]
TOPIC01_DIR = BASE_DIR / "专题产物" / "专题01"
TOPIC04_DIR = BASE_DIR / "专题产物" / "专题04"
SHEET4_EXCEL = TOPIC04_DIR / "表格" / "专题04_表格_四因素归因结果.xlsx"
OUT_PPT = TOPIC01_DIR / "汇报" / "专题01_汇报_归因瀑布图.pptx"
IMG_DIR = TOPIC04_DIR / "图表"
IMG_DIR.mkdir(parents=True, exist_ok=True)
MCK = {"dark_blue": "#051C2C", "light_blue": "#3A5F8A", "positive": "#27AE60", "negative": "#E74C3C", "warning": "#E67E22", "gray_dark": "#333333"}


def _slide_contains_text(slide, keyword):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if keyword in (r.text or ""):
                        return True
    return False


def remove_sheet4_slides_if_present(prs, max_remove=6):
    """若末尾连续页为专题四内容（含 Sheet④），则移除以便重新追加"""
    removed = 0
    while prs.slides and removed < max_remove:
        last = prs.slides[-1]
        if not _slide_contains_text(last, "Sheet④"):
            break
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])
        removed += 1


def draw_waterfall_contrib(filepath):
    """瀑布图：订单价、订单数、前台毛利率 对整体毛利额同比波动的贡献额及贡献占比"""
    tb = pd.read_excel(SHEET4_EXCEL, sheet_name="归因毛利额贡献")
    labels = tb["因素"].tolist()
    contrib_amt = tb["贡献额"].tolist()
    contrib_pct = tb["贡献占比_%"].tolist()
    n = len(labels)
    colors = [MCK["positive"] if v >= 0 else MCK["negative"] for v in contrib_amt]
    x = np.arange(n)
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.bar(x, contrib_amt, 0.5, color=colors)
    ax.axhline(0, color="gray", linewidth=0.8)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=11)
    ax.set_ylabel("对毛利额同比波动的贡献（元）", fontsize=11)
    ax.set_title("客品数、品单价、订单数、前台毛利率 对整体毛利额同比波动的贡献（YTD2026P1 vs 同期）", fontsize=12, pad=10)
    for i, (amt, pct) in enumerate(zip(contrib_amt, contrib_pct)):
        va = "bottom" if amt >= 0 else "top"
        ax.text(i, amt, f"{amt/1e6:.2f}M\n({pct:.1f}%)", ha="center", va=va, fontsize=9)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def add_slide_section(prs, title):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(1.2))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(5, 28, 44)
    tb.text_frame.paragraphs[0].alignment = 1


def add_slide_chart_simple(prs, title, img_path, footnote=""):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.55))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0, 0, 0)
    slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.0), width=Inches(12), height=Inches(5.5))
    if footnote:
        ft = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12), Inches(0.35))
        ft.text_frame.paragraphs[0].add_run().text = footnote
        ft.text_frame.paragraphs[0].font.size = Pt(9)
        ft.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)


def main():
    if not SHEET4_EXCEL.exists():
        print("请先运行 run_专题四_Sheet4_订单价订单数毛利率归因.py 生成", SHEET4_EXCEL)
        return
    p = IMG_DIR / "sheet4_waterfall_contrib.png"
    draw_waterfall_contrib(p)
    print("图表已生成:", p)
    if not OUT_PPT.exists():
        print("未找到", OUT_PPT)
        return
    prs = Presentation(str(OUT_PPT))
    remove_sheet4_slides_if_present(prs)
    add_slide_section(prs, "专题四：Sheet④ 客品数/品单价/订单数/前台毛利率 对毛利额同比归因")
    add_slide_chart_simple(
        prs,
        "客品数、品单价、订单数、前台毛利率 对整体毛利额同比波动的贡献（贡献额与贡献占比%）",
        p,
        "前台毛利额 = 客品数×品单价×订单数×前台毛利率，偏微分分解 dM≈P·N·G·dC+C·N·G·dP+C·P·G·dN+C·P·N·dG | 数据来源: 专题四 归因结果",
    )
    prs.save(OUT_PPT)
    print("已追加专题四 1 页至:", OUT_PPT)


if __name__ == "__main__":
    main()
