# -*- coding: utf-8 -*-
"""
专题三 Sheet③：生成 5 页图表并追加到 专题一_归因瀑布图.pptx
1. 帕累托图：TOP10 销售额贡献 SPU + 累计销售额占比，两时间窗口图例对比
2. 箱体图：X=箱体毛利率，图例=时间窗口，Y=SPU个数占比、毛利额占比
3. 复合图：X=箱体，图例=时间窗口，Y 左=SPU个数(柱)、Y 右=品效(线)、毛利额占比%(线)
4. 瀑布图：SPU 毛利率结构贡献（对整体毛利率的水平分解）
5. 瀑布图：SPU 毛利率波动贡献（对整体毛利率同比波动的归因）
"""
import os
from pathlib import Path
_BASE = Path(__file__).resolve().parent
os.environ["MPLCONFIGDIR"] = str(_BASE / ".mplcache")
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
plt.rcParams["font.sans-serif"] = ["PingFang SC", "Heiti SC", "STHeiti", "SimHei", "Microsoft YaHei", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE_DIR = Path(__file__).resolve().parents[2]
TOPIC01_DIR = BASE_DIR / "专题产物" / "专题01"
TOPIC03_DIR = BASE_DIR / "专题产物" / "专题03"
DATA_DIR = BASE_DIR / "原始数据"
SRC_EXCEL = DATA_DIR / "专题一：分析数据总表.xlsx"
SHEET3_EXCEL = TOPIC03_DIR / "表格" / "专题03_表格_SPU分析结果.xlsx"
OUT_PPT = TOPIC01_DIR / "汇报" / "专题01_汇报_归因瀑布图.pptx"
IMG_DIR = TOPIC03_DIR / "图表"
IMG_DIR.mkdir(parents=True, exist_ok=True)

MAT2026P1_MONTHS = [f"2025-{m:02d}" for m in range(2, 13)] + ["2026-01"]   # 2025-02～2026-01
MAT2025P1_MONTHS = [f"2024-{m:02d}" for m in range(2, 13)] + ["2025-01"]   # 2024-02～2025-01
MCK = {"dark_blue": "#051C2C", "light_blue": "#3A5F8A", "positive": "#27AE60", "negative": "#E74C3C", "warning": "#E67E22", "gray_dark": "#333333"}


def _slide_contains_text(slide, keyword):
    """检查幻灯片任意形状文本是否包含 keyword"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if keyword in (r.text or ""):
                        return True
    return False


def remove_sheet3_slides_if_present(prs, max_remove=12):
    """若末尾连续页为专题三内容（含 Sheet③），则移除以便重新追加（最多移除 max_remove 页）"""
    removed = 0
    while prs.slides and removed < max_remove:
        last = prs.slides[-1]
        if not _slide_contains_text(last, "Sheet③"):
            break
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])
        removed += 1


def draw_pareto_top10(filepath):
    """帕累托：TOP10 销售额 SPU（以 YTD2026P2 排序），两窗口销售额柱状 + 累计占比线"""
    df = pd.read_excel(SRC_EXCEL, sheet_name="③").dropna(subset=["SPU名称", "销售额"])
    def agg(w):
        return w.groupby("SPU名称")["销售额"].sum().sort_values(ascending=False).reset_index()
    w26 = df[df["月份"].isin(MAT2026P1_MONTHS)]
    w25 = df[df["月份"].isin(MAT2025P1_MONTHS)]
    s26 = agg(w26)
    s25 = agg(w25)
    s26["累计占比"] = s26["销售额"].cumsum() / s26["销售额"].sum()
    s25["累计占比"] = s25["销售额"].cumsum() / s25["销售额"].sum()
    top10_names = s26.head(10)["SPU名称"].tolist()
    short = [n[:12] + ".." if len(n) > 14 else n for n in top10_names]
    sale_26 = [s26[s26["SPU名称"] == n]["销售额"].sum() for n in top10_names]
    sale_25 = [s25[s25["SPU名称"] == n]["销售额"].sum() for n in top10_names]
    total26 = s26["销售额"].sum()
    total25 = s25["销售额"].sum()
    cum26, cum25 = [], []
    c26, c25 = 0, 0
    for i in range(10):
        c26 += sale_26[i]
        c25 += sale_25[i]
        cum26.append(c26 / total26 if total26 else 0)
        cum25.append(c25 / total25 if total25 else 0)
    x = np.arange(10)
    width = 0.35
    fig, ax1 = plt.subplots(figsize=(11, 5.5))
    ax1.bar(x - width/2, sale_26, width, label="MAT2026P1", color=MCK["dark_blue"])
    ax1.bar(x + width/2, sale_25, width, label="MAT2025P1", color=MCK["warning"])
    ax1.set_ylabel("销售额", fontsize=11)
    ax1.set_xlabel("SPU（按 MAT2026P1 销售额 TOP10）", fontsize=10)
    ax1.set_xticks(x)
    ax1.set_xticklabels(short, rotation=35, ha="right")
    ax1.legend(loc="upper left")
    ax2 = ax1.twinx()
    ax2.plot(x, [c*100 for c in cum26], "o-", color=MCK["dark_blue"], linewidth=2, markersize=6, label="MAT2026P1 累计占比%")
    ax2.plot(x, [c*100 for c in cum25], "s--", color=MCK["warning"], linewidth=2, markersize=6, label="MAT2025P1 累计占比%")
    ax2.set_ylabel("累计销售额占比（%）", fontsize=11)
    ax2.set_ylim(0, 105)
    ax2.legend(loc="upper right")
    ax1.set_title("TOP10 销售额贡献 SPU 及累计销售额占比 — 两时间窗口对比", fontsize=13, pad=10)
    fig.patch.set_facecolor("white")
    ax1.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def draw_bin_margin_chart(filepath):
    """毛利率等距 7 箱体：X=箱1～箱7（低毛利→高毛利），左轴=SPU个数占比，右轴=毛利额占比"""
    tb = pd.read_excel(SHEET3_EXCEL, sheet_name="数据表3_毛利率7分箱")
    order_col = "箱序" if "箱序" in tb.columns else "区间品线毛利率_中值"
    t26 = tb[tb["时间窗口"] == "MAT2026P1"].sort_values(order_col).reset_index(drop=True).head(7)
    t25 = tb[tb["时间窗口"] == "MAT2025P1"].sort_values(order_col).reset_index(drop=True).head(7)
    # 每窗口固定 7 行，按行取数并补足长度 7
    n26 = (t26["SPU非重复计数"].tolist() + [0] * 7)[:7]
    n25 = (t25["SPU非重复计数"].tolist() + [0] * 7)[:7]
    m26 = (t26["品线毛利额占比"].tolist() + [0.0] * 7)[:7]
    m25 = (t25["品线毛利额占比"].tolist() + [0.0] * 7)[:7]
    total_spu_26 = max(sum(n26), 1)
    total_spu_25 = max(sum(n25), 1)
    pct_spu_26 = [n / total_spu_26 * 100 for n in n26]
    pct_spu_25 = [n / total_spu_25 * 100 for n in n25]
    # X 轴：箱1=<0，箱2～7=0～最大毛利率等距
    labels = []
    for i in range(7):
        t = t26 if i < len(t26) else t25
        if i < len(t):
            cell = t.iloc[i].get("箱体", "")
            mid = t.iloc[i].get("区间品线毛利率_中值", 0)
            if cell == "<0":
                labels.append(f"箱{i+1}\n<0")
            else:
                labels.append(f"箱{i+1}\n({mid*100:.0f}%)")
        else:
            labels.append(f"箱{i+1}")
    x = np.arange(7)
    width = 0.36
    fig, ax1 = plt.subplots(figsize=(11, 5.5))
    ax1.bar(x - width/2, pct_spu_26, width, label="MAT2026P1 SPU个数占比", color=MCK["dark_blue"])
    ax1.bar(x + width/2, pct_spu_25, width, label="MAT2025P1 SPU个数占比", color=MCK["light_blue"], alpha=0.85)
    ax1.set_ylabel("SPU 个数占比（%）", fontsize=11)
    ax1.set_xlabel("毛利率分箱：箱1=负毛利(<0)，箱2～7=0～最大毛利率等距", fontsize=11)
    ax1.set_xticks(x)
    ax1.set_xticklabels(labels, fontsize=9)
    ax1.set_ylim(0, max(max(pct_spu_26, default=0), max(pct_spu_25, default=0)) * 1.15 or 100)
    ax1.legend(loc="upper right", fontsize=9)
    ax2 = ax1.twinx()
    ax2.plot(x, [a*100 for a in m26], "o-", color=MCK["dark_blue"], linewidth=2, markersize=7, label="MAT2026P1 毛利额占比")
    ax2.plot(x, [a*100 for a in m25], "s--", color=MCK["warning"], linewidth=2, markersize=6, label="MAT2025P1 毛利额占比")
    ax2.set_ylabel("品线毛利额占比（%）", fontsize=11)
    ax2.set_ylim(0, 105)
    ax2.legend(loc="upper left", fontsize=9)
    ax1.set_title("毛利率 7 箱体（<0 一箱 + 0～最大毛利率等距 6 箱）× SPU个数占比/毛利额占比", fontsize=12, pad=10)
    fig.patch.set_facecolor("white")
    ax1.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def draw_bin_composite_chart(filepath):
    """复合图：X=箱体，图例=时间窗口，Y 左=SPU个数(柱)，Y 右=品效(线)、毛利额占比%(线)"""
    tb = pd.read_excel(SHEET3_EXCEL, sheet_name="数据表3_毛利率7分箱")
    order_col = "箱序" if "箱序" in tb.columns else "区间品线毛利率_中值"
    t26 = tb[tb["时间窗口"] == "MAT2026P1"].sort_values(order_col).reset_index(drop=True).head(7)
    t25 = tb[tb["时间窗口"] == "MAT2025P1"].sort_values(order_col).reset_index(drop=True).head(7)
    n26 = (t26["SPU非重复计数"].tolist() + [0] * 7)[:7]
    n25 = (t25["SPU非重复计数"].tolist() + [0] * 7)[:7]
    eff26 = (t26["品效"].tolist() + [0.0] * 7)[:7]
    eff25 = (t25["品效"].tolist() + [0.0] * 7)[:7]
    m26 = (t26["品线毛利额占比"].tolist() + [0.0] * 7)[:7]
    m25 = (t25["品线毛利额占比"].tolist() + [0.0] * 7)[:7]
    # X 轴标签
    labels = []
    for i in range(7):
        t = t26 if i < len(t26) else t25
        if i < len(t):
            cell = t.iloc[i].get("箱体", "")
            labels.append(f"箱{i+1}\n<0" if cell == "<0" else f"箱{i+1}")
        else:
            labels.append(f"箱{i+1}")
    x = np.arange(7)
    width = 0.36
    fig, ax1 = plt.subplots(figsize=(11, 6))
    # 左轴：SPU 个数（柱状）
    b1 = ax1.bar(x - width/2, n26, width, label="MAT2026P1 SPU个数", color=MCK["dark_blue"])
    b2 = ax1.bar(x + width/2, n25, width, label="MAT2025P1 SPU个数", color=MCK["light_blue"], alpha=0.85)
    ax1.set_ylabel("SPU 个数", fontsize=11)
    ax1.set_xlabel("毛利率分箱（箱1=负毛利<0，箱2～7=0～最大毛利率等距）", fontsize=10)
    ax1.set_xticks(x)
    ax1.set_xticklabels(labels, fontsize=9)
    ax1.set_ylim(0, max(max(n26, default=0), max(n25, default=0)) * 1.2 or 1)
    ax1.legend(loc="upper right", fontsize=9)
    # 右轴：品效（线）
    ax2 = ax1.twinx()
    ax2.plot(x, eff26, "o-", color=MCK["dark_blue"], linewidth=2, markersize=7, label="MAT2026P1 品效")
    ax2.plot(x, eff25, "s--", color=MCK["warning"], linewidth=2, markersize=6, label="MAT2025P1 品效")
    ax2.set_ylabel("品效（毛利额/SPU）", fontsize=11)
    ymin_eff = min(min(eff26), min(eff25), 0)
    ax2.set_ylim(ymin_eff * 1.1 if ymin_eff < 0 else 0, max(max(eff26), max(eff25)) * 1.15 or 1)
    ax2.legend(loc="upper left", fontsize=9)
    # 右轴偏移：毛利额占比%（线）
    ax3 = ax1.twinx()
    ax3.spines["right"].set_position(("outward", 55))
    ax3.plot(x, [a*100 for a in m26], "^-", color=MCK["dark_blue"], linewidth=1.5, markersize=5, alpha=0.8, label="MAT2026P1 毛利额占比%")
    ax3.plot(x, [a*100 for a in m25], "v--", color=MCK["warning"], linewidth=1.5, markersize=5, alpha=0.8, label="MAT2025P1 毛利额占比%")
    ax3.set_ylabel("毛利额占比（%）", fontsize=11)
    ax3.set_ylim(0, max(max(m26)*100, max(m25)*100) * 1.15 or 100)
    ax3.legend(loc="center left", fontsize=9)
    ax1.set_title("箱体 × 时间窗口：SPU个数、品效、毛利额占比 复合图", fontsize=12, pad=10)
    fig.patch.set_facecolor("white")
    ax1.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def draw_waterfall_structure(filepath):
    """瀑布图：各 SPU 对整体毛利率的结构贡献（占比×毛利率），按贡献绝对值排序"""
    c = pd.read_excel(SHEET3_EXCEL, sheet_name="各SPU对毛利率波动贡献")
    total_sale_26 = c["销售额_26"].sum()
    c["结构贡献"] = (c["销售额_26"] / total_sale_26) * c["品线毛利率_26"]
    c = c.dropna(subset=["结构贡献"])
    c["abs"] = c["结构贡献"].abs()
    c = c.sort_values("abs", ascending=False).reset_index(drop=True)
    top_n = 12
    top = c.head(top_n)
    other_contrib = c["结构贡献"].sum() - top["结构贡献"].sum()
    overall_gp = c["结构贡献"].sum()
    labels = ["其他"] + top["SPU名称"].str[:10].tolist() + ["整体\n毛利率"]
    values = [other_contrib] + top["结构贡献"].tolist() + [overall_gp]
    colors = [MCK["gray_dark"]] + [MCK["positive"] if v >= 0 else MCK["negative"] for v in top["结构贡献"]] + [MCK["dark_blue"]]
    n = len(labels)
    b2, h2 = [0], [values[0]]
    for i in range(1, n - 1):
        b2.append(b2[-1] + h2[-1])
        h2.append(values[i])
    b2.append(0)
    h2.append(values[-1])
    fig, ax = plt.subplots(figsize=(11, 5.5))
    for i in range(n):
        if i == 0 or i == n - 1:
            ax.bar(i, h2[i], 0.5, bottom=0, color=colors[i])
        else:
            ax.bar(i, h2[i], 0.5, bottom=b2[i], color=colors[i])
    ax.set_xticks(range(n))
    ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=9)
    ax.set_ylabel("对整体毛利率的结构贡献", fontsize=11)
    ax.set_title("SPU 对整体毛利率的结构贡献（按贡献绝对值 TOP12 + 其他）", fontsize=13, pad=10)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f"{x*100:.1f}%"))
    for i in range(n):
        if i == 0 or i == n - 1:
            ax.text(i, h2[i]/2, f"{h2[i]*100:.1f}%", ha="center", va="center", fontsize=8, color="white")
        else:
            ax.text(i, b2[i] + h2[i]/2, f"{h2[i]*100:+.1f}%", ha="center", va="center", fontsize=7, color="white")
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def draw_waterfall_change(filepath):
    """瀑布图：各 SPU 对整体毛利率同比波动的贡献（贡献率 pp），按绝对值排序；Y 轴截断使 TOP12 柱更直观"""
    c = pd.read_excel(SHEET3_EXCEL, sheet_name="各SPU对毛利率波动贡献")
    c = c.dropna(subset=["贡献率_pp"])
    c["abs"] = c["贡献率_pp"].abs()
    c = c.sort_values("abs", ascending=False).reset_index(drop=True)
    top_n = 12
    top = c.head(top_n)
    other_contrib = c["贡献率_pp"].sum() - top["贡献率_pp"].sum()
    delta_gp = c["贡献率_pp"].sum()
    labels = ["其他"] + top["SPU名称"].str[:10].tolist() + ["整体\n同比波动"]
    values = [other_contrib] + top["贡献率_pp"].tolist() + [delta_gp]
    colors = [MCK["gray_dark"]] + [MCK["positive"] if v >= 0 else MCK["negative"] for v in top["贡献率_pp"]] + [MCK["dark_blue"]]
    n = len(labels)
    b2, h2 = [0], [values[0]]
    for i in range(1, n - 1):
        b2.append(b2[-1] + h2[-1])
        h2.append(values[i])
    b2.append(0)
    h2.append(values[-1])
    # 计算除第一根外柱子的 Y 范围，用于截断 Y 轴
    rest_bottoms = b2[1:]
    rest_tops = [b + hv for b, hv in zip(rest_bottoms, h2[1:])]
    y_rest_min = min(min(rest_bottoms), min(rest_tops), 0, delta_gp)
    y_rest_max = max(max(rest_bottoms), max(rest_tops), 0, delta_gp)
    margin = max((y_rest_max - y_rest_min) * 0.15, 0.3)
    y_lim_low = y_rest_min - margin
    y_lim_high = y_rest_max + margin
    # 若第一根「其他」柱绝对值远大于其余柱范围，则截断 Y 轴
    need_break = abs(other_contrib) > (y_lim_high - y_lim_low) * 0.6
    if need_break and other_contrib > y_lim_high:
        # 第一根柱只画到 y_lim_high，柱顶标注真实值
        first_bar_height_visible = y_lim_high
    elif need_break and other_contrib < y_lim_low:
        first_bar_height_visible = y_lim_low - other_contrib  # 负柱可见高度
    else:
        first_bar_height_visible = None

    fig, ax = plt.subplots(figsize=(11, 5.5))
    for i in range(n):
        if i == 0:
            if first_bar_height_visible is not None and other_contrib > y_lim_high:
                ax.bar(i, first_bar_height_visible, 0.5, bottom=0, color=colors[i])
                ax.text(i, y_lim_high + (y_lim_high - y_lim_low) * 0.04, f"{h2[i]:.2f}pp", ha="center", va="bottom", fontsize=8, color=MCK["gray_dark"])
            elif first_bar_height_visible is not None and other_contrib < y_lim_low:
                ax.bar(i, first_bar_height_visible, 0.5, bottom=other_contrib, color=colors[i])
                ax.text(i, y_lim_low - (y_lim_high - y_lim_low) * 0.04, f"{h2[i]:.2f}pp", ha="center", va="top", fontsize=8, color=MCK["gray_dark"])
            else:
                ax.bar(i, h2[i], 0.5, bottom=0, color=colors[i])
                ax.text(i, h2[i]/2, f"{h2[i]:.2f}pp", ha="center", va="center", fontsize=8, color="white")
        elif i == n - 1:
            ax.bar(i, h2[i], 0.5, bottom=0, color=colors[i])
            ax.text(i, h2[i]/2, f"{h2[i]:.2f}pp", ha="center", va="center", fontsize=8, color="white")
        else:
            ax.bar(i, h2[i], 0.5, bottom=b2[i], color=colors[i])
            ax.text(i, b2[i] + h2[i]/2, f"{h2[i]:+.2f}pp", ha="center", va="center", fontsize=7, color="white")
    y_show_high = y_lim_high + (y_lim_high - y_lim_low) * 0.12 if need_break and other_contrib > y_lim_high else y_lim_high
    ax.set_ylim(y_lim_low, y_show_high)
    ax.set_xticks(range(n))
    ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=9)
    ax.set_ylabel("对整体毛利率同比波动的贡献（pp）" + (" · Y轴已截断" if need_break else ""), fontsize=11)
    ax.set_title("SPU 对整体毛利率同比波动的归因贡献（按绝对值 TOP12 + 其他）", fontsize=13, pad=10)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def add_slide_section(prs, title):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(1.2))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(5, 28, 44)
    tb.text_frame.paragraphs[0].alignment = 1


def add_slide_chart_simple(prs, title, img_path, footnote=""):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.55))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0, 0, 0)
    slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.0), width=Inches(12), height=Inches(5.5))
    if footnote:
        ft = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12), Inches(0.35))
        ft.text_frame.paragraphs[0].add_run().text = footnote
        ft.text_frame.paragraphs[0].font.size = Pt(9)
        ft.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)


def main():
    if not SHEET3_EXCEL.exists():
        print("请先运行 run_专题三_Sheet3_SPU分析.py 生成", SHEET3_EXCEL)
        return
    p1 = IMG_DIR / "sheet3_pareto_top10.png"
    p2 = IMG_DIR / "sheet3_bin_margin.png"
    p3 = IMG_DIR / "sheet3_bin_composite.png"
    p4 = IMG_DIR / "sheet3_waterfall_structure.png"
    p5 = IMG_DIR / "sheet3_waterfall_change.png"
    draw_pareto_top10(p1)
    draw_bin_margin_chart(p2)
    draw_bin_composite_chart(p3)
    draw_waterfall_structure(p4)
    draw_waterfall_change(p5)
    print("图表已生成:", p1, p2, p3, p4, p5)
    if not OUT_PPT.exists():
        print("未找到", OUT_PPT)
        return
    prs = Presentation(str(OUT_PPT))
    remove_sheet3_slides_if_present(prs)
    add_slide_section(prs, "专题三：Sheet③ SPU 分析")
    add_slide_chart_simple(prs, "TOP10 销售额贡献 SPU 及累计销售额占比 — MAT2026P1 / MAT2025P1 对比", p1, "数据来源: 专题一分析数据总表 Sheet③")
    add_slide_chart_simple(prs, "箱体毛利率 × SPU个数占比 / 毛利额占比 — 两时间窗口", p2, "数据来源: 专题三_Sheet3_SPU分析结果 数据表3")
    add_slide_chart_simple(prs, "箱体 × 时间窗口：SPU个数、品效、毛利额占比 复合图", p3, "左轴=SPU个数(柱)，右轴=品效(线)、毛利额占比%(线) | 数据来源: 专题三 数据表3")
    add_slide_chart_simple(prs, "SPU 对整体毛利率的结构贡献（水平分解，按贡献绝对值 TOP12）", p4, "结构贡献 = 销售额占比 × 该SPU毛利率 | 数据来源: 专题三")
    add_slide_chart_simple(prs, "SPU 对整体毛利率同比波动的归因贡献（按绝对值 TOP12）", p5, "贡献率 = 销售额占比 × (当期毛利率 - 同期毛利率) | 数据来源: 专题三")
    prs.save(OUT_PPT)
    print("已追加专题三 5 页图表至:", OUT_PPT)


if __name__ == "__main__":
    main()
