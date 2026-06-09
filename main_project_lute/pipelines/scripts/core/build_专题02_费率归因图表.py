# MOVED TO: scripts/phase1/build_专题02_费率归因图表.py
# -*- coding: utf-8 -*-
"""
专题二 Sheet②：各费率对品线毛利率同比波动的贡献占比图，并追加到 专题一_归因瀑布图.pptx
依赖：专题二_Sheet2_双平台费率归因结果.xlsx 已生成（先运行 run_专题二_Sheet2_双平台费率归因.py）
"""
import os
from pathlib import Path
_BASE = Path(__file__).resolve().parent
os.environ["MPLCONFIGDIR"] = str(_BASE / ".mplcache")
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
plt.rcParams["font.sans-serif"] = ["PingFang SC", "Heiti SC", "STHeiti", "SimHei", "Microsoft YaHei", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE_DIR = Path(__file__).resolve().parents[2]
TOPIC01_DIR = BASE_DIR / "专题产物" / "专题01"
TOPIC02_DIR = BASE_DIR / "专题产物" / "专题02"
EXCEL_PATH = TOPIC02_DIR / "表格" / "专题02_表格_双平台费率归因结果.xlsx"
OUT_PPT = TOPIC01_DIR / "汇报" / "专题01_汇报_归因瀑布图.pptx"
IMG_DIR = TOPIC02_DIR / "图表"
IMG_DIR.mkdir(parents=True, exist_ok=True)

MCK = {
    "dark_blue": "#051C2C",
    "light_blue": "#3A5F8A",
    "positive": "#27AE60",
    "negative": "#E74C3C",
    "warning": "#E67E22",
    "gray_dark": "#333333",
}

FEE_LABELS = {
    "退款率": "退款率",
    "促销折扣率": "促销折扣",
    "线上推广费率(品线)": "线上推广",
    "生产成本率": "生产成本",
    "头程费率(营收)": "头程",
    "仓储配送费率": "仓储配送",
    "平台佣金率": "平台佣金",
    "其他费率": "其他",
}
RATE_COLS = list(FEE_LABELS.keys())


def _slide_contains_text(slide, keyword):
    """检查幻灯片任意形状文本是否包含 keyword"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if keyword in (r.text or ""):
                        return True
    return False


def remove_sheet2_slides_if_present(prs, max_remove=8):
    """若末尾连续页为专题二内容（含 Sheet②），则移除以便重新追加（最多移除 max_remove 页）"""
    removed = 0
    while prs.slides and removed < max_remove:
        last = prs.slides[-1]
        if not _slide_contains_text(last, "Sheet②"):
            break
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])
        removed += 1


def draw_fee_structure_waterfall(tab, period_label, filepath):
    """
    各费率对品线毛利率的结构贡献（占比）瀑布图：营收100% → 各费率扣减 → 品线毛利率。
    一图两子图：左 亚马逊，右 独立站；图例为独立站、亚马逊（通过子图标题区分）。
    """
    # 费率顺序（与 RATE_COLS 一致），瀑布从 100% 依次扣减
    rate_order = RATE_COLS  # 退款率, 促销折扣率, ..., 其他费率
    short_labels = [FEE_LABELS[c] for c in rate_order]
    x_labels = ["营收\n100%"] + short_labels + ["品线\n毛利率"]

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5.5))
    for ax, platform, color, title in [
        (ax1, "亚马逊", MCK["dark_blue"], "亚马逊"),
        (ax2, "独立站", MCK["warning"], "独立站"),
    ]:
        row = tab[tab["平台"] == platform].iloc[0]
        gp = float(row.get("品线毛利率_当期", 0) or 0)
        rates = [float(row.get(f"{c}_当期", 0) or 0) for c in rate_order]
        # 使用绝对值（表中费率可能为负表示成本）
        rates = [abs(r) for r in rates]
        # 瀑布：首柱 100%，中间各柱为扣减段（高度=该费率），末柱=品线毛利率
        n = len(rates) + 2
        bottoms = [0.0]
        heights = [1.0]
        cum = 1.0
        for r in rates:
            cum -= r
            bottoms.append(cum)
            heights.append(r)
        bottoms.append(0.0)
        heights.append(gp)

        colors = [MCK["dark_blue"] if platform == "亚马逊" else MCK["warning"]]
        for i in range(len(rates)):
            colors.append(MCK["light_blue"] if platform == "亚马逊" else MCK["gray_dark"])
        colors.append(MCK["positive"])

        x = np.arange(n)
        bar_w = 0.5
        for i in range(n):
            if i == 0 or i == n - 1:
                ax.bar(i, heights[i], bar_w, bottom=0, color=colors[i], edgecolor="none")
            else:
                ax.bar(i, heights[i], bar_w, bottom=bottoms[i], color=colors[i], edgecolor="none")

        ax.set_xticks(x)
        ax.set_xticklabels(x_labels, rotation=30, ha="right", fontsize=9)
        ax.set_ylabel("占比", fontsize=10)
        ax.set_title(title, fontsize=12, fontweight="bold")
        ax.set_ylim(0, 1.05)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f"{x*100:.0f}%"))
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        # 柱上标注占比值
        for i in range(n):
            if i == 0:
                ax.text(i, heights[i] / 2, "100%", ha="center", va="center", fontsize=8, color="white", fontweight="bold")
            elif i == n - 1:
                ax.text(i, heights[i] / 2, f"{heights[i]*100:.1f}%", ha="center", va="center", fontsize=8, color="white", fontweight="bold")
            else:
                mid = bottoms[i] + heights[i] / 2
                ax.text(i, mid, f"{heights[i]*100:.1f}%", ha="center", va="center", fontsize=7, color="white")
    fig.suptitle(f"各费率对品线毛利率的结构贡献（占比）— {period_label}", fontsize=13, fontweight="bold", y=1.02)
    fig.patch.set_facecolor("white")
    ax1.set_facecolor("white")
    ax2.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()
    return filepath


def draw_fee_contribution_chart(tab, period_label, filepath):
    """
    绘制各费率对品线毛利率同比波动的贡献占比（%），亚马逊 vs 独立站 分组柱状图。
    柱上标注：贡献占比数值 + 正贡献/负贡献。
    """
    am = tab[tab["平台"] == "亚马逊"].iloc[0]
    id_ = tab[tab["平台"] == "独立站"].iloc[0]

    labels = [FEE_LABELS[c] for c in RATE_COLS]
    x = np.arange(len(labels))
    width = 0.35
    am_vals = [am.get(f"{c}_对毛利率波动贡献占比%", np.nan) for c in RATE_COLS]
    id_vals = [id_.get(f"{c}_对毛利率波动贡献占比%", np.nan) for c in RATE_COLS]
    am_vals = [v if pd.notna(v) else 0 for v in am_vals]
    id_vals = [v if pd.notna(v) else 0 for v in id_vals]

    fig, ax = plt.subplots(figsize=(11, 6))
    bars1 = ax.bar(x - width / 2, am_vals, width, label="亚马逊", color=MCK["dark_blue"])
    bars2 = ax.bar(x + width / 2, id_vals, width, label="独立站", color=MCK["warning"])
    ax.axhline(0, color="gray", linestyle="--", linewidth=0.8)

    # 柱上数据标注：贡献占比数值 + 正贡献/负贡献（按 y 轴范围做小幅偏移）
    y_min, y_max = ax.get_ylim()
    y_span = y_max - y_min if y_max > y_min else 20
    offset = max(5, y_span * 0.03)
    for bar, val in zip(bars1, am_vals):
        if pd.notna(val) and abs(val) >= 0.05:
            y_pos = bar.get_height() + (offset if val >= 0 else -offset)
            va = "bottom" if val >= 0 else "top"
            label = f"{val:+.1f}%\n{'正贡献' if val > 0 else '负贡献'}"
            ax.text(bar.get_x() + bar.get_width() / 2, y_pos, label, ha="center", va=va, fontsize=7, color=MCK["dark_blue"])
    for bar, val in zip(bars2, id_vals):
        if pd.notna(val) and abs(val) >= 0.05:
            y_pos = bar.get_height() + (offset if val >= 0 else -offset)
            va = "bottom" if val >= 0 else "top"
            label = f"{val:+.1f}%\n{'正贡献' if val > 0 else '负贡献'}"
            ax.text(bar.get_x() + bar.get_width() / 2, y_pos, label, ha="center", va=va, fontsize=7, color=MCK["warning"])

    ax.set_ylabel("贡献占比（%）", fontsize=11)
    ax.set_title(f"各费率对品线毛利率同比波动的贡献占比 — {period_label}", fontsize=13, pad=10)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, rotation=25, ha="right")
    ax.legend(loc="upper right", fontsize=9)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()
    return filepath


def add_slide_section(prs, section_title):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(1.2))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = section_title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(5, 28, 44)
    tb.text_frame.paragraphs[0].alignment = 1
    return slide


def add_slide_chart_simple(prs, title, img_path, footnote=""):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.55))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0, 0, 0)
    slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.0), width=Inches(12), height=Inches(5.5))
    if footnote:
        ft = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12), Inches(0.35))
        ft.text_frame.paragraphs[0].add_run().text = footnote
        ft.text_frame.paragraphs[0].font.size = Pt(9)
        ft.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
    return slide


def add_slide_conclusions(prs, lines):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(6.8))
    tb.text_frame.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tb.text_frame.paragraphs[0]
        else:
            p = tb.text_frame.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(14) if (line.startswith("一、") or line.startswith("二、") or line.startswith("三、")) else Pt(11)
        r.font.bold = line.startswith("一、") or line.startswith("二、") or line.startswith("三、")
        r.font.color.rgb = RGBColor(0, 0, 0) if r.font.bold else RGBColor(51, 51, 51)
    return slide


def main():
    if not EXCEL_PATH.exists():
        print("请先运行 run_专题二_Sheet2_双平台费率归因.py 生成", EXCEL_PATH)
        return
    tab_mat = pd.read_excel(EXCEL_PATH, sheet_name="MAT2026P1_vs_MAT2025P1")
    tab_ytd = pd.read_excel(EXCEL_PATH, sheet_name="YTD2026P1_vs_YTD2025P1")

    path_mat = IMG_DIR / "sheet2_各费率对毛利率波动贡献占比_MAT.png"
    path_ytd = IMG_DIR / "sheet2_各费率对毛利率波动贡献占比_YTD.png"
    path_struct_mat = IMG_DIR / "sheet2_费率结构瀑布图_MAT.png"
    path_struct_ytd = IMG_DIR / "sheet2_费率结构瀑布图_YTD.png"
    draw_fee_contribution_chart(tab_mat, "MAT2026P1 vs MAT2025P1", path_mat)
    draw_fee_contribution_chart(tab_ytd, "YTD2026P1 vs YTD2025P1", path_ytd)
    draw_fee_structure_waterfall(tab_mat, "MAT2026P1（2025-02～2026-01）", path_struct_mat)
    draw_fee_structure_waterfall(tab_ytd, "YTD2026P1（2026年1月）", path_struct_ytd)
    print("图表已生成:", path_mat, path_ytd, path_struct_mat, path_struct_ytd)

    if not OUT_PPT.exists():
        print("未找到", OUT_PPT, "，请先运行 build_waterfall_ppt.py 生成 PPT。")
        return
    prs = Presentation(str(OUT_PPT))
    remove_sheet2_slides_if_present(prs)

    add_slide_section(prs, "专题二：Sheet② 亚马逊 vs 独立站 费率归因")
    add_slide_chart_simple(
        prs,
        "各费率对品线毛利率同比波动的贡献占比 — MAT2026P1 vs MAT2025P1",
        path_mat,
        "贡献占比 = 该费率贡献率(pp) / 品线毛利率同比差值(pp)×100 | 数据来源: 专题一分析数据总表 Sheet②",
    )
    add_slide_chart_simple(
        prs,
        "各费率对品线毛利率同比波动的贡献占比 — YTD2026P1 vs YTD2025P1",
        path_ytd,
        "数据来源: 专题一分析数据总表 Sheet②",
    )
    add_slide_chart_simple(
        prs,
        "各费率对品线毛利率的结构贡献（占比）— MAT2026P1（图例：亚马逊、独立站）",
        path_struct_mat,
        "营收100% 依次扣减各费率，得到品线毛利率 | 数据来源: 专题一分析数据总表 Sheet②",
    )
    add_slide_chart_simple(
        prs,
        "各费率对品线毛利率的结构贡献（占比）— YTD2026P1（图例：亚马逊、独立站）",
        path_struct_ytd,
        "营收100% 依次扣减各费率，得到品线毛利率 | 数据来源: 专题一分析数据总表 Sheet②",
    )
    insight_lines = [
        "一、公式（Momcozy 管理层口径）",
        "• 卖出一块钱：硬性成本分两类 → 前台成本（促销+推广+退款）+ 后台成本（生产+头程+仓配+佣金+其他），剩下来的才是毛利。",
        "• 口径说明：前/后台占比（如 27%/73%）是“总成本内部结构”；折算到每1元收入需再乘以总成本率。",
        "",
        "二、卖出一块钱，钱花哪儿了？（当期结构）",
        "• 亚马逊：前台/后台在总成本内约 27%/73%；折算到每1元收入约 2毛1/5毛6，毛利约 2毛3。",
        "• 独立站：前台/后台在总成本内约 51%/49%；折算到每1元收入约 4毛7~4毛9/4毛6，毛利仅 6~7分。",
        "",
        "三、今年比去年：谁在变好、谁在变差？",
        "• 亚马逊：毛利率同比升（MAT +3.4pp，YTD +2.0pp），主要来自前台成本率下降（少打折/少推广/退款控好）。",
        "• 独立站：毛利率同比降（MAT −2.5pp，YTD −6.1pp），YTD 恶化主因是前台大幅多花（约 −8.4pp 贡献），推广/促销/退款压力大。",
        "",
        "四、关键词",
        "Momcozy | 前台成本 | 后台成本 | 品线毛利率 | 数据来源: 专题一分析数据总表 Sheet②",
    ]
    add_slide_conclusions(prs, insight_lines)
    business_lines = [
        "五、前台成本 — 结论与行动方向（Momcozy 吸奶器等母婴用品）",
        "• 含义：促销折扣 + 线上推广 + 退款（顾客与流量侧）。",
        "• 亚马逊：前台在总成本内约 27%（折算收入约 21%~22%）且同比在降 → 保持节奏，不盲目加促销/推广。",
        "• 独立站：前台在总成本内约 51%（折算收入约 47%~49%）且 YTD 恶化 → 优先 ① 控推广 ② 控促销 ③ 控退款。",
        "",
        "六、后台成本 — 结论与行动方向",
        "• 含义：生产成本 + 头程 + 仓配 + 平台佣金 + 其他（供应链与履约侧）。",
        "• 亚马逊：后台略升，持续盯采购与头程；FBA 与结构优化可巩固。",
        "• 独立站：无佣金但生产/头程/仓配仍有压力；先稳住前台，再系统压后台（供应商、头程方式、关税与海外仓）。",
        "",
        "七、给领导的三条结论",
        "① 独立站卖1块钱约有4毛7~4毛9花在前台，毛利仅6~7分；要改善毛利必须先压前台。② 今年独立站恶化主因是前台多花（推广/促销/退款），优先控推广、控促销、控退款。③ 后台在稳住前台后再系统优化。详见 专题二_前台与后台成本洞察_Momcozy.txt",
    ]
    add_slide_conclusions(prs, business_lines)
    prs.save(OUT_PPT)
    print("已追加专题二 Sheet② 费率归因图表与结论页至:", OUT_PPT)


if __name__ == "__main__":
    main()
