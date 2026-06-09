# -*- coding: utf-8 -*-
"""
供应链六大成本 36% 分析与优化 — 汇报 PPT 生成

面向 IT 总监 & 供应链总监，从数据分析维度推进供应链优化项目
15-20 页，麦肯锡风格
"""

import os
from pathlib import Path
import numpy as np

_BASE = Path(__file__).resolve().parent
os.environ["MPLCONFIGDIR"] = str(_BASE / ".mplcache")

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
plt.rcParams["font.sans-serif"] = ["PingFang SC", "Heiti SC", "STHeiti", "SimHei", "Microsoft YaHei"]
plt.rcParams["axes.unicode_minus"] = False

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = PROJECT_ROOT / "ref" / "books" / "供应链36%方案_Page6节点"
OUT_PPT = OUT_DIR / "供应链六大成本36%分析_IT与供应链总监汇报.pptx"
CHARTS_DIR = OUT_DIR / "_charts"
CHARTS_DIR.mkdir(parents=True, exist_ok=True)

MCK = {
    "dark_blue": "#051C2C",
    "medium_blue": "#1E3A5F",
    "light_blue": "#3A5F8A",
    "pale_blue": "#5A7F9F",
    "positive": "#27AE60",
    "negative": "#E74C3C",
    "orange": "#FF6B35",
}


def hex_to_rgb(hex_color):
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(MCK["dark_blue"])


def add_content_slide(prs, title, lines, font_size=18):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(MCK["dark_blue"])
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.space_after = Pt(6)
        p.font.color.rgb = hex_to_rgb("#333333")


def add_chart_slide(prs, title, chart_path, bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = Pt(28)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(MCK["dark_blue"])
    if chart_path and Path(chart_path).exists():
        slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1), width=Inches(9))
    if bullets:
        desc = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.8))
        tf = desc.text_frame
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(11)
            p.font.color.rgb = hex_to_rgb(MCK["medium_blue"])


def gen_cost_structure_chart():
    nodes = ["采购", "头程", "仓储", "尾程", "退换补发", "小包直邮"]
    pct = [20, 8, 6, 7, 4, 3]
    colors = [MCK["dark_blue"], MCK["medium_blue"], MCK["light_blue"],
              MCK["pale_blue"], MCK["orange"], MCK["negative"]]
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.barh(nodes, pct, color=colors)
    ax.axvline(x=36, color=MCK["negative"], linestyle="--", linewidth=2, label="当前 36%")
    ax.axvline(x=28, color=MCK["positive"], linestyle="--", linewidth=1.5, label="目标 28%")
    ax.set_xlabel("占销售额 %")
    ax.set_title("六大成本分项占比（示意，待补齐取数口径）")
    ax.legend(loc="lower right")
    ax.set_xlim(0, 42)
    for i, (n, v) in enumerate(zip(nodes, pct)):
        ax.text(v + 0.3, i, f"{v}%", va="center", fontsize=11)
    plt.tight_layout()
    path = CHARTS_DIR / "cost_structure.png"
    plt.savefig(path, dpi=150)
    plt.close()
    return path


def gen_node_matrix_chart():
    nodes = ["①采购", "②头程", "③仓储", "④尾程", "⑤退换补发", "⑥小包直邮"]
    p0 = ["防杂减重 降本三台阶", "方式优化 批量频次", "治乱 周转提升", "方式优化", "退货率 部分退", "占比优化"]
    fig, ax = plt.subplots(figsize=(10, 5))
    y = np.arange(len(nodes))
    ax.barh(y, [1] * len(nodes), color=[MCK["light_blue"]] * len(nodes), height=0.6)
    ax.set_yticks(y)
    ax.set_yticklabels(nodes, fontsize=12)
    ax.set_xticks([])
    for i, t in enumerate(p0):
        ax.text(0.5, i, t, ha="center", va="center", fontsize=10)
    ax.set_title("六大节点 P0 优化方向")
    plt.tight_layout()
    path = CHARTS_DIR / "node_matrix.png"
    plt.savefig(path, dpi=150)
    plt.close()
    return path


def gen_analytics_view_chart():
    views = ["成本结构看板", "成本-时效矩阵", "供应商绩效", "库存健康度", "退货归因", "部分退组合"]
    fig, ax = plt.subplots(figsize=(10, 5))
    x = np.arange(len(views))
    colors = [MCK["dark_blue"], MCK["medium_blue"], MCK["light_blue"],
             MCK["pale_blue"], MCK["orange"], MCK["negative"]]
    ax.bar(x, [1]*6, color=colors)
    ax.set_xticks(x)
    ax.set_xticklabels(views, rotation=15, ha="right")
    ax.set_title("分析视图与大数据算法选型")
    plt.tight_layout()
    path = CHARTS_DIR / "analytics_view.png"
    plt.savefig(path, dpi=150)
    plt.close()
    return path


def gen_it_supply_chain_chart():
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.axis("off")
    boxes = [
        ("IT 能力建设", 0.2, 0.6, 0.35, 0.25),
        ("供应链业务", 0.55, 0.6, 0.35, 0.25),
        ("数据取数口径", 0.2, 0.25, 0.35, 0.2),
        ("分析视图落地", 0.55, 0.25, 0.35, 0.2),
    ]
    for text, x, y, w, h in boxes:
        rect = plt.Rectangle((x, y), w, h, fill=True, facecolor=MCK["light_blue"], edgecolor=MCK["dark_blue"])
        ax.add_patch(rect)
        ax.text(x + w/2, y + h/2, text, ha="center", va="center", fontsize=11)
    ax.text(0.5, 0.05, "共同目标：36% -> 28%-30%", ha="center", fontsize=12, fontweight="bold")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    path = CHARTS_DIR / "it_supply_chain.png"
    plt.savefig(path, dpi=150)
    plt.close()
    return path


def main():
    print("=== 供应链六大成本 36% 分析 — IT与供应链总监汇报 PPT ===\n")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "供应链六大成本 36% 分析与优化",
        "从数据分析维度推进供应链优化项目 | 汇报对象：IT 总监 & 供应链总监")

    add_content_slide(prs, "目录", [
        "1. 问题与目标", "2. 36% 成本框架与六大节点",
        "3. 按节点：战略–战术–战斗 + 指标 + 分析视图 + 优化",
        "4. 数据分析能力建设（IT 视角）", "5. 分析视图与算法选型",
        "6. 取数口径与分项占比", "7. 行动计划与双部门协作",
    ], font_size=20)

    add_content_slide(prs, "一、问题与目标", [
        "【现状】供应链成本占销售额 36%，高于母婴跨境健康区间（25%–30%）",
        "六大成本：采购、头程、仓储、尾程、退/换/补发、小包直邮",
        "", "【目标】建立成本–节点统一框架，目标降至 28%–30%",
        "【汇报价值】以数据分析为牵引，推动 IT 与供应链协同落地优化项目",
    ])

    add_content_slide(prs, "二、36% 成本框架与六大节点", [
        "【口径】分子=六大成本合计 / 分母=销售额",
        "① 采购(15%-25%) | ② 头程(5%-12%) | ③ 仓储(3%-8%)",
        "④ 尾程(5%-10%) | ⑤ 退换补发(2%-6%) | ⑥ 小包直邮",
    ])

    chart1 = gen_cost_structure_chart()
    add_chart_slide(prs, "六大成本分项占比（示意）", chart1, ["待建立取数口径，优先实现「供应链成本结构看板」"])

    add_content_slide(prs, "三、六大节点总览", [
        "① 采购：防杂减重、降本三台阶 | ② 头程：方式优化、批量频次",
        "③ 仓储：治乱、周转提升 | ④ 尾程：方式优化、时效成本平衡",
        "⑤ 退/换/补发：退货率下降、部分退优化、VOC 闭环",
        "⑥ 小包直邮：占比优化、线路伙伴",
    ])

    chart2 = gen_node_matrix_chart()
    add_chart_slide(prs, "六大节点 P0 优化方向", chart2)

    add_content_slide(prs, "节点① 采购 — 战略战术战斗", [
        "【战略】防杂减重、降本三台阶；采购参与需求管理",
        "【指标】采购成本率、供应商准时率、质量合格率、集中采购占比",
        "【分析视图】采购成本结构看板、供应商绩效矩阵、补货优化",
        "【P0 优化】防杂减重、料号→品类→供应商系统降本",
    ])

    add_content_slide(prs, "节点②③ 头程与仓储", [
        "【头程】方式优化（海运/空运/铁路）、批量与频次、清关合规",
        "【仓储】治乱（预测+安全库存+补货）、周转提升、仓网优化",
        "【分析视图】头程成本–时效矩阵；库存健康度、周转呆滞预警",
    ])

    add_content_slide(prs, "节点④⑤⑥ 尾程、退换补发、小包直邮", [
        "【尾程】方式优化、时效成本平衡",
        "【退/换/补发】退货率下降、部分退优化、退款→VOC 闭环",
        "【小包直邮】占比优化、直邮 vs 海外仓对比",
    ])

    add_content_slide(prs, "四、数据分析能力建设（IT 视角）", [
        "【数据基础】建立供应链成本主题域，与 01 矩阵、05 数仓表对齐",
        "【分析视图】优先：供应链成本结构看板、库存健康度、退货归因",
        "【算法】因素分解、瀑布图、安全库存公式、分类/帕累托、关联规则",
    ])

    chart3 = gen_analytics_view_chart()
    add_chart_slide(prs, "分析视图与大数据算法选型", chart3, [
        "成本结构→因素分解/瀑布图 | 库存健康度→安全库存/ABC | 退货归因→分类/帕累托",
    ])

    add_content_slide(prs, "五、取数口径与分项占比", [
        "采购成本：采购系统、成本分摊表 | 头程：物流对账单",
        "仓储：海外仓/FBA 费用报表 | 尾程：配送费、平台扣费",
        "退换补发：退货处理费、换货/补发物流 | 小包直邮：直邮物流对账单",
        "", "【建议】在 01 矩阵或 05 数仓中增加「供应链成本结构」主题",
    ])

    chart4 = gen_it_supply_chain_chart()
    add_chart_slide(prs, "六、IT 与供应链双部门协作", chart4)

    add_content_slide(prs, "七、行动计划与里程碑", [
        "【Phase 1】补齐分项占比（1–2 月）IT：数据主题建模；供应链：业务口径确认",
        "【Phase 2】落地分析视图（2–3 月）IT：看板开发；供应链：指标与阈值确认",
        "【Phase 3】P0 优化落地（持续）双部门联合复盘，迭代指标体系",
    ])

    add_content_slide(prs, "下一步行动", [
        "1. 补齐分项占比 2. 落地分析视图 3. P0 优化落地 4. 双部门协作",
        "", "【交付物】ref/books/供应链36%方案_Page6节点/",
    ])

    add_content_slide(prs, "附录：参考书籍（books 20–31）", [
        "20 宫迅伟 | 21 采购与供应链大数据 | 22 赵先德 | 23 马向国",
        "25 佟昕 | 26 刘宝红·三道防线 | 29 刘宝红·高成本高库存 | 31 刘宝红·采购",
    ], font_size=14)

    prs.save(OUT_PPT)
    print(f"PPT 已生成: {OUT_PPT}")
    print(f"共 {len(prs.slides)} 页")


if __name__ == "__main__":
    main()
