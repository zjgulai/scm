# MOVED TO: scripts/phase3/build_phase3_topic3_ppt.py
# -*- coding: utf-8 -*-
"""
Phase3 专题③ 渠道健康度 PPT
"""

import os
from pathlib import Path

_BASE = Path(__file__).resolve().parent
os.environ["MPLCONFIGDIR"] = str(_BASE / ".mplcache")

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
plt.rcParams["font.sans-serif"] = ["PingFang SC", "Heiti SC"]

from pptx import Presentation
from pptx.util import Pt

PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUTPUT_DIR = PROJECT_ROOT / "phase3_outputs" / "topic3_channel"
OUT_PPT = PROJECT_ROOT / "phase3_outputs" / "Topic3_Channel_汇报.pptx"


def add_slide(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    body.text_frame.clear()
    for line in lines:
        p = body.text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)


def main():
    prs = Presentation()

    # 标题页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "专题③ 渠道健康度分析"
    slide.placeholders[1].text = "分国家渠道运营健康度提升"

    # 执行摘要
    add_slide(prs, "执行摘要", [
        "✓ 渠道生命周期分析：识别各渠道所处阶段",
        "✓ 渠道差异性分析：自然/付费/红人流量结构",
        "✓ 风险预警：高风险与高机会渠道识别",
    ])

    # 生命周期
    lifecycle = OUTPUT_DIR / "channel_lifecycle.csv"
    if lifecycle.exists():
        import pandas as pd
        df = pd.read_csv(lifecycle)
        lines = ["渠道生命周期："]
        for _, row in df.iterrows():
            lines.append(f"• {row['channel_id']}: {row['lifecycle_stage']}")
        add_slide(prs, "一、渠道生命周期", lines)

    # 风险机会
    risk = OUTPUT_DIR / "channel_risk_opportunity.csv"
    if risk.exists():
        df = pd.read_csv(risk)
        high_risk = df[df["category"] == "高风险"]
        high_opp = df[df["category"] == "高机会"]

        lines = ["风险与机会："]
        if not high_risk.empty:
            lines.append("\n⚠️ 高风险：")
            for _, row in high_risk.head(3).iterrows():
                lines.append(f"  {row['channel_id']}/{row['country_code']}")
        if not high_opp.empty:
            lines.append("\n✅ 高机会：")
            for _, row in high_opp.head(3).iterrows():
                lines.append(f"  {row['channel_id']}/{row['country_code']}")
        add_slide(prs, "二、风险与机会", lines)

    # 结论
    add_slide(prs, "三、驱动动作", [
        "1. 对高风险渠道制定改进计划",
        "2. 加大高机会渠道资源投入",
        "3. 优化流量结构配比",
        "4. 调整库存策略",
    ])

    prs.save(OUT_PPT)
    print(f"PPT已生成: {OUT_PPT}")


if __name__ == "__main__":
    main()
