# -*- coding: utf-8 -*-
"""
Phase3 专题④ 营销ROI PPT
"""

import os
from pathlib import Path

_BASE = Path(__file__).resolve().parent
os.environ["MPLCONFIGDIR"] = str(_BASE / ".mplcache")

from pptx import Presentation
from pptx.util import Pt

PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUTPUT_DIR = PROJECT_ROOT / "phase3_outputs" / "topic4_marketing"
OUT_PPT = PROJECT_ROOT / "phase3_outputs" / "Topic4_Marketing_汇报.pptx"


def add_slide(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    body.text_frame.clear()
    for line in lines:
        p = body.text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)


def main():
    prs = Presentation()

    # 标题页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "专题④ 营销ROI分析"
    slide.placeholders[1].text = "营销项目类型ROI量化提升"

    # 执行摘要
    add_slide(prs, "执行摘要", [
        "✓ 营销活动效果分析：各活动类型ROI对比",
        "✓ 费用结构分析：CPM/CPC效率指标",
        "✓ 活动ROI排名：高ROI/低ROI活动识别",
    ])

    # 活动效果
    lifecycle = OUTPUT_DIR / "campaign_lifecycle_analysis.csv"
    if lifecycle.exists():
        import pandas as pd
        df = pd.read_csv(lifecycle)
        lines = ["各活动类型ROI："]
        for _, row in df.iterrows():
            lines.append(f"• {row['campaign_type']}: ROI={row['roi']:.2f}, CTR={row['ctr']:.2f}%")
        add_slide(prs, "一、营销活动效果", lines)

    # 费用效率
    cost = OUTPUT_DIR / "cost_structure.csv"
    if cost.exists():
        df = pd.read_csv(cost)
        latest = df[df["dt_month"] == df["dt_month"].max()]
        if not latest.empty:
            row = latest.iloc[0]
            lines = ["最新月份费用效率：", f"• 广告花费: ${row['ad_spend']:,.0f}", f"• CPM: ${row['cpm']:.2f}", f"• CPC: ${row['cpc']:.2f}"]
            add_slide(prs, "二、费用效率", lines)

    # ROI排名
    roi = OUTPUT_DIR / "campaign_roi_by_spu.csv"
    if roi.exists():
        df = pd.read_csv(roi)
        high = df[df["performance"] == "高ROI"]
        lines = ["高ROI活动："]
        if not high.empty:
            for _, row in high.head(5).iterrows():
                lines.append(f"• {row['campaign_id']}: ROI={row['roi']:.2f}")
        add_slide(prs, "三、ROI排名", lines)

    # 结论
    add_slide(prs, "四、驱动动作", [
        "1. 加大高ROI活动投入",
        "2. 优化低ROI活动预算",
        "3. 降低CPM/CPC成本",
        "4. 建立ROI监控机制",
    ])

    prs.save(OUT_PPT)
    print(f"PPT已生成: {OUT_PPT}")


if __name__ == "__main__":
    main()
