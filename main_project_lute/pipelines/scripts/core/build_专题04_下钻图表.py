# MOVED TO: scripts/phase1/build_专题04_下钻图表.py
# -*- coding: utf-8 -*-
"""
专题四 Sheet④ 下钻：客品数归因 TOP10、毛利率归因 TOP10 瀑布图，追加到 专题一_归因瀑布图.pptx
"""
import os
from pathlib import Path
_BASE = Path(__file__).resolve().parent
os.environ["MPLCONFIGDIR"] = str(_BASE / ".mplcache")
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
plt.rcParams["font.sans-serif"] = ["PingFang SC", "Heiti SC", "STHeiti", "SimHei", "Microsoft YaHei", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE_DIR = Path(__file__).resolve().parents[2]
TOPIC01_DIR = BASE_DIR / "专题产物" / "专题01"
TOPIC04_DIR = BASE_DIR / "专题产物" / "专题04"
DRILLDOWN_EXCEL = TOPIC04_DIR / "表格" / "专题04_表格_下钻归因SPU明细.xlsx"
TOP10_C = TOPIC04_DIR / "表格" / "专题04_表格_客品数归因TOP10.xlsx"
TOP10_G = TOPIC04_DIR / "表格" / "专题04_表格_毛利率归因TOP10.xlsx"
OUT_PPT = TOPIC01_DIR / "汇报" / "专题01_汇报_归因瀑布图.pptx"
IMG_DIR = TOPIC04_DIR / "图表"
IMG_DIR.mkdir(parents=True, exist_ok=True)
MCK = {"dark_blue": "#051C2C", "light_blue": "#3A5F8A", "positive": "#27AE60", "negative": "#E74C3C", "warning": "#E67E22", "gray_dark": "#333333"}


def _slide_contains_text(slide, keyword):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if keyword in (r.text or ""):
                        return True
    return False


def remove_drilldown_slides_if_present(prs, max_remove=4):
    """若末尾连续页为专题四下钻图表（整体客品数归因/整体毛利率归因），则移除"""
    removed = 0
    while prs.slides and removed < max_remove:
        last = prs.slides[-1]
        if not (_slide_contains_text(last, "整体客品数归因") or _slide_contains_text(last, "整体毛利率归因")):
            break
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])
        removed += 1


def draw_waterfall_kepinshu(filepath):
    """客品数同比波动 SPU 归因 TOP10 瀑布图"""
    full = pd.read_excel(DRILLDOWN_EXCEL)
    top10 = pd.read_excel(TOP10_C)
    delta_c = full["客品数_贡献"].sum()
    other = delta_c - top10["客品数_贡献"].sum()
    # 横轴：仅一个「其余SPU」+ TOP10 SPU 名称 + 一个「整体变化」，避免出现两个「其他」
    short_names = [str(x)[:10] if str(x) != "其他" else "其余" for x in top10["SPU名称"].tolist()]
    labels = ["其余SPU"] + short_names + ["整体\n变化"]
    values = [other] + top10["客品数_贡献"].tolist() + [delta_c]
    colors = [MCK["gray_dark"]] + [MCK["positive"] if v >= 0 else MCK["negative"] for v in top10["客品数_贡献"]] + [MCK["dark_blue"]]
    n = len(labels)
    b2, h2 = [0], [values[0]]
    for i in range(1, n - 1):
        b2.append(b2[-1] + h2[-1])
        h2.append(values[i])
    b2.append(0)
    h2.append(values[-1])
    fig, ax = plt.subplots(figsize=(11, 5.5))
    for i in range(n):
        if i == 0 or i == n - 1:
            ax.bar(i, h2[i], 0.5, bottom=0, color=colors[i])
        else:
            ax.bar(i, h2[i], 0.5, bottom=b2[i], color=colors[i])
    ax.axhline(0, color="gray", linewidth=0.6)
    ax.set_xticks(range(n))
    ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=9)
    ax.set_ylabel("对整体客品数同比波动的贡献", fontsize=11)
    ax.set_title("整体客品数归因 — 贡献绝对值 TOP10 SPU（YTD2026P2 vs YTD2025P2 同期）", fontsize=12, pad=10)
    for i in range(n):
        if i == 0 or i == n - 1:
            ax.text(i, h2[i]/2, f"{h2[i]:.4f}", ha="center", va="center", fontsize=8, color="white")
        else:
            ax.text(i, b2[i] + h2[i]/2, f"{h2[i]:.4f}", ha="center", va="center", fontsize=7, color="white")
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def draw_waterfall_maolilv(filepath):
    """毛利率同比波动 SPU 归因 TOP10 瀑布图"""
    full = pd.read_excel(DRILLDOWN_EXCEL)
    top10 = pd.read_excel(TOP10_G)
    delta_gp = full["毛利率_贡献"].sum()
    other = delta_gp - top10["毛利率_贡献"].sum()
    short_names = [str(x)[:10] if str(x) != "其他" else "其余" for x in top10["SPU名称"].tolist()]
    labels = ["其余SPU"] + short_names + ["整体\n变化"]
    values = [other] + top10["毛利率_贡献"].tolist() + [delta_gp]
    values_pp = [v * 100 for v in values]
    colors = [MCK["gray_dark"]] + [MCK["positive"] if v >= 0 else MCK["negative"] for v in top10["毛利率_贡献"]] + [MCK["dark_blue"]]
    n = len(labels)
    b2, h2 = [0], [values_pp[0]]
    for i in range(1, n - 1):
        b2.append(b2[-1] + h2[-1])
        h2.append(values_pp[i])
    b2.append(0)
    h2.append(values_pp[-1])
    fig, ax = plt.subplots(figsize=(11, 5.5))
    for i in range(n):
        if i == 0 or i == n - 1:
            ax.bar(i, h2[i], 0.5, bottom=0, color=colors[i])
        else:
            ax.bar(i, h2[i], 0.5, bottom=b2[i], color=colors[i])
    ax.axhline(0, color="gray", linewidth=0.6)
    ax.set_xticks(range(n))
    ax.set_xticklabels(labels, rotation=30, ha="right", fontsize=9)
    ax.set_ylabel("对整体毛利率同比波动的贡献（pp）", fontsize=11)
    ax.set_title("整体毛利率归因 — 贡献绝对值 TOP10 SPU（YTD2026P2 vs YTD2025P2 同期）", fontsize=12, pad=10)
    for i in range(n):
        if i == 0 or i == n - 1:
            ax.text(i, h2[i]/2, f"{h2[i]:.2f}pp", ha="center", va="center", fontsize=8, color="white")
        else:
            ax.text(i, b2[i] + h2[i]/2, f"{h2[i]:.2f}pp", ha="center", va="center", fontsize=7, color="white")
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def add_slide_chart_simple(prs, title, img_path, footnote=""):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.55))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0, 0, 0)
    slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.0), width=Inches(12), height=Inches(5.5))
    if footnote:
        ft = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12), Inches(0.35))
        ft.text_frame.paragraphs[0].add_run().text = footnote
        ft.text_frame.paragraphs[0].font.size = Pt(9)
        ft.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)


def main():
    if not TOP10_C.exists() or not TOP10_G.exists():
        print("请先运行 run_专题四_Sheet4_客品数与毛利率下钻归因.py 生成 TOP10 与明细")
        return
    p1 = IMG_DIR / "sheet4_drilldown_kepinshu.png"
    p2 = IMG_DIR / "sheet4_drilldown_maolilv.png"
    draw_waterfall_kepinshu(p1)
    draw_waterfall_maolilv(p2)
    print("图表已生成:", p1, p2)
    if not OUT_PPT.exists():
        print("未找到", OUT_PPT)
        return
    prs = Presentation(str(OUT_PPT))
    remove_drilldown_slides_if_present(prs)
    add_slide_chart_simple(prs, "整体客品数归因 — 贡献绝对值 TOP10 SPU（YTD2026P2 vs YTD2025P2 同期）", p1,
        "整体客品数 = Σ(订单数_s/Σ订单数)×客品数_s | 时间: YTD2026P2 / YTD2025P2 | 数据来源: 专题四 下钻归因")
    add_slide_chart_simple(prs, "整体毛利率归因 — 贡献绝对值 TOP10 SPU（YTD2026P2 vs YTD2025P2 同期）", p2,
        "整体毛利率 = Σ(销售额_s/Σ销售额)×毛利率_s | 时间: YTD2026P2 / YTD2025P2 | 数据来源: 专题四 下钻归因")
    prs.save(OUT_PPT)
    print("已追加专题四下钻 2 页至:", OUT_PPT)


if __name__ == "__main__":
    main()
