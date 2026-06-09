# -*- coding: utf-8 -*-
"""
PPT Generator - PPT生成模块（重构版）

变更：
- 使用 python-pptx 生成真正的 .pptx 文件
- 支持标题页、洞察页、图表页（插入PNG图片）
- 保留 JSON 输出作为调试接口
"""

import json
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUTPUT_DIR = PROJECT_ROOT / "tmp" / "outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# 麦肯锡品牌色
MCKINSEY_BLUE = RGBColor(0x05, 0x1C, 0x2C)
MCKINSEY_ACCENT = RGBColor(0xFF, 0x6B, 0x35)
MCKINSEY_LIGHT_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)


class PPTGenerator:
    """PPT 生成器 — 使用 python-pptx"""

    def __init__(self):
        self.output_dir = OUTPUT_DIR

    def generate(self, analysis_result: Dict, output_name: str = None) -> Path:
        """根据分析结果生成 PPTX 文件"""
        skill_name = analysis_result.get("skill_name", "report")
        title = analysis_result.get("title", "分析报告")
        insights = analysis_result.get("insights", [])
        chart_images = analysis_result.get("chart_images", [])
        summary = analysis_result.get("summary", {})
        next_steps = analysis_result.get("next_steps", [])

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 第 1 页：标题
        self._add_title_slide(prs, title, skill_name)

        # 第 2 页：核心洞察
        if insights:
            self._add_insights_slide(prs, "核心发现", insights)

        # 第 3-N 页：数据摘要
        if summary:
            self._add_summary_slide(prs, "数据摘要", summary)

        # 图表页：插入 PNG 图片
        for img_path in chart_images:
            if Path(img_path).exists():
                self._add_image_slide(prs, "图表分析", img_path)

        # 建议措施页
        if next_steps:
            self._add_insights_slide(prs, "建议措施", next_steps)

        # 保存
        output_path = self.output_dir / f"{output_name or skill_name}.pptx"
        prs.save(str(output_path))
        print(f"[PPTGenerator] 已生成: {output_path}")
        return output_path

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str = ""):
        """麦肯锡风格标题页"""
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # 蓝色背景条
        from pptx.util import Inches, Pt
        left, top, width, height = Inches(0), Inches(0), prs.slide_width, Inches(2.5)
        shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = MCKINSEY_BLUE
        shape.line.fill.background()

        # 标题文字
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 副标题
        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = f"卓越商业分析专家AI SaaS | {subtitle}"
            p2.font.size = Pt(14)
            p2.font.color.rgb = MCKINSEY_LIGHT_BLUE

        # 日期
        txBox3 = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        from datetime import datetime
        p3.text = f"生成时间: {datetime.now().strftime('%Y-%m-%d')}"
        p3.font.size = Pt(10)
        p3.font.color.rgb = LIGHT_GREY

    def _add_insights_slide(self, prs: Presentation, title: str, items: List[str]):
        """麦肯锡风格洞察页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 顶部蓝色条
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = MCKINSEY_BLUE
        shape.line.fill.background()

        # 标题
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 洞察要点（每行一个）
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(12)

    def _add_summary_slide(self, prs: Presentation, title: str, summary: Dict):
        """数据摘要页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 顶部栏
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = MCKINSEY_BLUE
        shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 数据摘要
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        lines = self._flatten_summary(summary)
        for i, line in enumerate(lines):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(8)

    def _add_image_slide(self, prs: Presentation, title: str, img_path: str):
        """图表页 — 插入 PNG 图片"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 标题
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = MCKINSEY_BLUE
        shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 插入图片（居中）
        try:
            slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.2), Inches(10), Inches(5.5))
        except Exception as e:
            txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(10), Inches(1))
            p2 = txBox2.text_frame.paragraphs[0]
            p2.text = f"[图表加载失败: {img_path}]"
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)

    def _flatten_summary(self, summary: Dict) -> List[str]:
        """将嵌套 summary 展平为字符串列表"""
        lines = []
        for key, value in summary.items():
            if isinstance(value, dict):
                for k, v in value.items():
                    lines.append(f"{k}: {v}")
            elif isinstance(value, list):
                lines.append(f"{key}: {len(value)}项")
            elif isinstance(value, float):
                lines.append(f"{key}: {value:.2%}" if abs(value) < 10 else f"{key}: {value:,.2f}")
            else:
                lines.append(f"{key}: {value}")
        return lines

    def generate_from_result(self, result, output_name: str = None) -> Path:
        """从 AnalysisResult 对象生成 PPT"""
        result_dict = {
            "skill_name": getattr(result, "skill_name", "report"),
            "title": getattr(result, "title", "分析报告"),
            "summary": getattr(result, "summary", {}),
            "insights": getattr(result, "insights", []),
            "chart_images": getattr(result, "chart_images", []),
            "next_steps": getattr(result, "next_steps", []),
        }
        return self.generate(result_dict, output_name)

    def generate_json(self, analysis_result: Dict, output_name: str = None) -> Path:
        """调试用 JSON 输出"""
        output_path = self.output_dir / f"{output_name or 'debug'}.json"
        output_path.write_text(json.dumps(analysis_result, ensure_ascii=False, indent=2))
        print(f"[PPTGenerator] JSON: {output_path}")
        return output_path
