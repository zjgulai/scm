# MOVED TO: scripts/phase3/build_phase3_topic1_ppt.py
# -*- coding: utf-8 -*-
"""
Phase3 专题① VOC汇报PPT生成
"""

import os
from pathlib import Path
import pandas as pd

_BASE = Path(__file__).resolve().parent
os.environ["MPLCONFIGDIR"] = str(_BASE / ".mplcache")

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
plt.rcParams["font.sans-serif"] = ["PingFang SC", "Heiti SC", "STHeiti"]
plt.rcParams["axes.unicode_minus"] = False

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUTPUT_DIR = PROJECT_ROOT / "phase3_outputs" / "topic1_voc"
OUT_PPT = PROJECT_ROOT / "phase3_outputs" / "Topic1_VOC_汇报.pptx"

MCK = {
    "dark_blue": "#051C2C",
    "medium_blue": "#1E3A5F",
    "light_blue": "#3A5F8A",
    "positive": "#27AE60",
    "negative": "#E74C3C",
}


def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_content_slide(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    body.text_frame.clear()
    for line in lines:
        p = body.text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)


def generate_charts():
    """生成分析图表"""
    # 货架内痛点
    df = pd.read_csv(OUTPUT_DIR / "shelf_inside_analysis.csv")
    top_issues = df.nlargest(10, "priority")

    fig, ax = plt.subplots(figsize=(10, 6))
    colors = [MCK["negative"] if r > 0.15 else MCK["light_blue"] for r in top_issues["bad_rate"]]
    ax.barh(top_issues["spu_id"] + " " + top_issues["country_code"], top_issues["bad_rate"] * 100, color=colors)
    ax.set_xlabel("差评率 (%)")
    ax.set_title("货架内Top痛点SPU")
    plt.tight_layout()
    chart1 = OUTPUT_DIR / "chart_pain_points.png"
    plt.savefig(chart1, dpi=150)
    plt.close()

    # 品牌提及
    df2 = pd.read_csv(OUTPUT_DIR / "shelf_outside_brand_mention.csv")
    df2 = df2.sort_values("post_cnt", ascending=True)

    fig, ax = plt.subplots(figsize=(10, 6))
    colors = [MCK["positive"] if b == "Momcozy" else MCK["light_blue"] for b in df2["brand"]]
    ax.barh(df2["brand"], df2["post_cnt"], color=colors)
    ax.set_xlabel("提及次数")
    ax.set_title("社媒品牌提及排名")
    plt.tight_layout()
    chart2 = OUTPUT_DIR / "chart_brand_mention.png"
    plt.savefig(chart2, dpi=150)
    plt.close()

    return chart1, chart2


def load_insights():
    txt_path = OUTPUT_DIR / "topic1_insights.txt"
    if txt_path.exists():
        with open(txt_path, "r", encoding="utf-8") as f:
            return f.read()
    return ""


def main():
    print("=== 生成 Topic1 VOC 汇报 PPT ===\n")

    prs = Presentation()

    # 标题页
    add_title_slide(prs, "专题① VOC数据洞察", "全域用户声音分析 - 货架内/外/竞品")

    # 执行摘要
    add_content_slide(prs, "执行摘要", [
        "✓ 货架内分析：识别产品痛点与亮点",
        "✓ 货架外分析：发现高潜需求与第二大单品机会",
        "✓ 竞品分析：提炼本土化营销话术",
        "✓ 趋势分析：指导渠道布局与流量投入",
    ])

    # 货架内
    add_content_slide(prs, "一、货架内用户声音", [
        "核心发现：",
        "• SPU209在AU差评率18.8%",
        "• SPU202在UK差评率20.8%",
        "• DTC渠道差评率整体高于AMZ",
        "",
        "驱动动作：优化产品细节，提升用户满意度",
    ])

    chart1, _ = generate_charts()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text_frame.paragraphs[0].text = "货架内Top痛点"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    slide.shapes.add_picture(str(chart1), Inches(0.5), Inches(1), width=Inches(9))

    # 货架外
    add_content_slide(prs, "二、货架外高潜需求", [
        "核心发现：",
        "• 高浓度社区：BabyCenter, Reddit, Mumsnet",
        "• Momcozy社媒提及284次，好评率49.3%",
        "• Spectra/Medela竞品好评率更高",
        "",
        "驱动动作：在高浓度社区增加内容投放",
    ])

    # 竞品
    competitor = pd.read_csv(OUTPUT_DIR / "competitor_opportunity.csv")
    if not competitor.empty:
        opp_text = ["竞品机会分析："]
        for _, row in competitor.head(3).iterrows():
            opp_text.append(f"• {row['country']}/{row['channel']}: 比{row['competitor']}高{row['advantage']:.1f}分")
    else:
        opp_text = ["暂无竞品机会数据"]
    add_content_slide(prs, "三、竞品机会", opp_text)

    # 结论
    insights = load_insights()
    add_content_slide(prs, "四、核心结论", insights.split("\n")[:15])

    # 保存
    prs.save(OUT_PPT)
    print(f"PPT已生成: {OUT_PPT}")


if __name__ == "__main__":
    main()
