# -*- coding: utf-8 -*-
"""
共享 PPT 配置 — python-pptx 公共工具函数

所有 build_*_ppt.py 脚本通过此模块获取统一的 PPT 样式。
"""

from pathlib import Path
from datetime import datetime
from typing import List, Dict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

PROJECT_ROOT = Path(__file__).resolve().parents[4]
PPT_OUTPUT_DIR = PROJECT_ROOT / "tmp" / "outputs"
PPT_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# 麦肯锡品牌色
MCKINSEY_BLUE = RGBColor(0x05, 0x1C, 0x2C)
MCKINSEY_ACCENT = RGBColor(0xFF, 0x6B, 0x35)
MCKINSEY_LIGHT_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def create_presentation() -> Presentation:
    """创建标准 16:9 演示文稿"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    """添加麦肯锡风格标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MCKINSEY_BLUE
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(1.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = f"卓越商业分析专家AI SaaS | {subtitle}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = MCKINSEY_LIGHT_BLUE

    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
    p3 = txBox3.text_frame.paragraphs[0]
    p3.text = f"生成时间: {datetime.now().strftime('%Y-%m-%d')}"
    p3.font.size = Pt(10)
    p3.font.color.rgb = LIGHT_GREY


def add_section_header(prs: Presentation, title: str):
    """添加章节标题条"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MCKINSEY_BLUE
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return slide


def add_bullet_slide(prs: Presentation, section_title: str, items: List[str]):
    """添加要点列表页"""
    slide = add_section_header(prs, section_title)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(12)
    return slide


def add_image_slide(prs: Presentation, title: str, img_path: str):
    """添加图表页 — 居中插入 PNG"""
    slide = add_section_header(prs, title)
    try:
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.2), Inches(10), Inches(5.5))
    except Exception:
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(10), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"[图表: {img_path}]"
        p.font.size = Pt(14)
        p.font.color.rgb = MCKINSEY_ACCENT
    return slide


def save_presentation(prs: Presentation, filename: str, output_dir: Path = None) -> Path:
    """统一 PPT 保存"""
    out_dir = output_dir or PPT_OUTPUT_DIR
    out_dir.mkdir(parents=True, exist_ok=True)
    if not filename.endswith('.pptx'):
        filename += '.pptx'
    path = out_dir / filename
    prs.save(str(path))
    print(f"[ppt_common] 已保存: {path}")
    return path
