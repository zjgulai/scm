# -*- coding: utf-8 -*-
"""
专题五：购物篮关联分析 ⑤⑥⑦⑧ — 制作对比分析图表并追加到 专题一_归因瀑布图.pptx
"""
import os
from pathlib import Path
_BASE = Path(__file__).resolve().parent
os.environ["MPLCONFIGDIR"] = str(_BASE / ".mplcache")
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
plt.rcParams["font.sans-serif"] = ["PingFang SC", "Heiti SC", "STHeiti", "SimHei", "Microsoft YaHei", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE_DIR = Path(__file__).resolve().parents[2]
TOPIC01_DIR = BASE_DIR / "专题产物" / "专题01"
TOPIC05_DIR = BASE_DIR / "专题产物" / "专题05"
BASKET_EXCEL = TOPIC05_DIR / "表格" / "专题05_表格_购物篮关联分析汇总.xlsx"
OUT_PPT = TOPIC01_DIR / "汇报" / "专题01_汇报_归因瀑布图.pptx"
IMG_DIR = TOPIC05_DIR / "图表"
IMG_DIR.mkdir(parents=True, exist_ok=True)
MCK = {"dark_blue": "#051C2C", "light_blue": "#3A5F8A", "positive": "#27AE60", "negative": "#E74C3C", "warning": "#E67E22", "gray_dark": "#333333"}


def _slide_contains_text(slide, keyword):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if keyword in (r.text or ""):
                        return True
    return False


def remove_sheet5_slides_if_present(prs, max_remove=15):
    removed = 0
    while prs.slides and removed < max_remove:
        last = prs.slides[-1]
        if not (_slide_contains_text(last, "购物篮") or _slide_contains_text(last, "专题五")):
            break
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])
        removed += 1


def draw_product_line_top_lift(filepath):
    """⑤ 产品线层级：提升度 TOP10 规则（先继项→后继项）"""
    df = pd.read_excel(BASKET_EXCEL, sheet_name="⑤_提升度≥1")
    df = df.nlargest(10, "提升度")
    df["标签"] = df["先继项_product_line"] + "→" + df["后继项_product_line"]
    labels = df["标签"].tolist()
    lift = df["提升度"].tolist()
    colors = [MCK["positive"] if v >= 1 else MCK["negative"] for v in lift]
    x = np.arange(len(labels))
    fig, ax = plt.subplots(figsize=(10, 5.5))
    ax.barh(x, lift, color=colors)
    ax.set_yticks(x)
    ax.set_yticklabels(labels, fontsize=10)
    ax.set_xlabel("提升度", fontsize=11)
    ax.axvline(1, color="gray", linestyle="--", linewidth=1)
    ax.set_title("⑤ 订单-产品线 关联规则 TOP10（按提升度，提升度≥1）", fontsize=12, pad=10)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def draw_lift_ge1_vs_lt1(filepath):
    """⑤ 提升度≥1 与 提升度<1 规则数量对比及占比"""
    d5_ge1 = pd.read_excel(BASKET_EXCEL, sheet_name="⑤_提升度≥1")
    d5_lt1 = pd.read_excel(BASKET_EXCEL, sheet_name="⑤_提升度<1")
    n_ge1, n_lt1 = len(d5_ge1), len(d5_lt1)
    fig, ax = plt.subplots(figsize=(6, 4))
    ax.bar([0], [n_ge1], 0.5, label="提升度≥1（正关联）", color=MCK["positive"])
    ax.bar([1], [n_lt1], 0.5, label="提升度<1（互斥）", color=MCK["negative"])
    ax.set_xticks([0, 1])
    ax.set_xticklabels(["提升度≥1", "提升度<1"])
    ax.set_ylabel("规则数", fontsize=11)
    ax.set_title("⑤ 订单-产品线 规则分布：正关联 vs 互斥", fontsize=12, pad=10)
    for i, v in enumerate([n_ge1, n_lt1]):
        ax.text(i, v + 0.5, str(v), ha="center", va="bottom", fontsize=11)
    ax.legend()
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def draw_spu_top_lift(filepath):
    """⑥ 订单-SPU 提升度 TOP10（先继项→后继项，用 SPU 名称简称）"""
    df = pd.read_excel(BASKET_EXCEL, sheet_name="⑥_订单SPU")
    df = df.nlargest(10, "提升度")
    labels = [f"{a[:8]}→{b[:8]}" for a, b in zip(df["先继项_spu_name"], df["后继项_spu_name"])]
    lift = df["提升度"].tolist()
    x = np.arange(len(labels))
    fig, ax = plt.subplots(figsize=(10, 5.5))
    ax.barh(x, lift, color=MCK["light_blue"])
    ax.set_yticks(x)
    ax.set_yticklabels(labels, fontsize=9)
    ax.set_xlabel("提升度", fontsize=11)
    ax.set_title("⑥ 订单-SPU 关联规则 TOP10（按提升度）", fontsize=12, pad=10)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def _pivot_lift_for_heatmap(df, name_col_ante, name_col_post, lift_col="提升度", top_n=12):
    """取先继项/后继项出现频次 TOP 的实体，构造 先继项×后继项 提升度矩阵。"""
    from collections import Counter
    all_names = list(df[name_col_ante].dropna()) + list(df[name_col_post].dropna())
    top_names = [x for x, _ in Counter(all_names).most_common(top_n)]
    if len(top_names) < 2:
        return None, None, None
    # 只保留 先继项、后继项 均在 top_names 中的规则，取 (先,后) 对应提升度（若有多条取最大）
    sub = df[df[name_col_ante].isin(top_names) & df[name_col_post].isin(top_names)]
    pivot = sub.groupby([name_col_ante, name_col_post])[lift_col].max().unstack(fill_value=0)
    # 对齐行列到 top_names，缺失为 0
    for n in top_names:
        if n not in pivot.index:
            pivot.loc[n] = 0
        if n not in pivot.columns:
            pivot[n] = 0
    pivot = pivot.reindex(index=top_names, columns=top_names, fill_value=0)
    return pivot, top_names, top_names


def draw_heatmap_customer_spu(filepath):
    """⑦ 客户-SPU：先继项×后继项 提升度热力图（TOP SPU）"""
    df = pd.read_excel(BASKET_EXCEL, sheet_name="⑦_客户SPU")
    pivot, rows, cols = _pivot_lift_for_heatmap(
        df, "先继项_spu_name", "后继项_spu_name", top_n=12
    )
    if pivot is None or pivot.size == 0:
        plt.figure(figsize=(6, 5))
        plt.text(0.5, 0.5, "⑦ 数据不足", ha="center", va="center", fontsize=14)
        plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
        plt.close()
        return
    labels = [str(x)[:10] + "…" if len(str(x)) > 10 else str(x) for x in rows]
    fig, ax = plt.subplots(figsize=(9, 7))
    im = ax.imshow(pivot.values, cmap="YlOrRd", aspect="auto", vmin=0, vmax=max(1, pivot.values.max()))
    ax.set_xticks(np.arange(len(cols)))
    ax.set_yticks(np.arange(len(rows)))
    ax.set_xticklabels(labels, fontsize=8, rotation=45, ha="right")
    ax.set_yticklabels(labels, fontsize=8)
    for i in range(len(rows)):
        for j in range(len(cols)):
            v = pivot.iloc[i, j]
            ax.text(j, i, f"{v:.2f}" if v > 0 else "", ha="center", va="center", fontsize=7)
    ax.set_title("⑦ 客户-SPU 关联提升度热力图（先继项→后继项，TOP SPU）", fontsize=12, pad=10)
    plt.colorbar(im, ax=ax, label="提升度")
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def draw_heatmap_customer_sku(filepath):
    """⑧ 客户-SKU：先继项×后继项 提升度热力图（TOP SKU）"""
    df = pd.read_excel(BASKET_EXCEL, sheet_name="⑧_客户SKU")
    pivot, rows, cols = _pivot_lift_for_heatmap(
        df, "先继项_sku_name", "后继项_sku_name", top_n=12
    )
    if pivot is None or pivot.size == 0:
        plt.figure(figsize=(6, 5))
        plt.text(0.5, 0.5, "⑧ 数据不足", ha="center", va="center", fontsize=14)
        plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
        plt.close()
        return
    labels = [str(x)[:8] + "…" if len(str(x)) > 8 else str(x) for x in rows]
    fig, ax = plt.subplots(figsize=(9, 7))
    im = ax.imshow(pivot.values, cmap="YlOrRd", aspect="auto", vmin=0, vmax=max(1, pivot.values.max()))
    ax.set_xticks(np.arange(len(cols)))
    ax.set_yticks(np.arange(len(rows)))
    ax.set_xticklabels(labels, fontsize=8, rotation=45, ha="right")
    ax.set_yticklabels(labels, fontsize=8)
    for i in range(len(rows)):
        for j in range(len(cols)):
            v = pivot.iloc[i, j]
            ax.text(j, i, f"{v:.2f}" if v > 0 else "", ha="center", va="center", fontsize=7)
    ax.set_title("⑧ 客户-SKU 关联提升度热力图（先继项→后继项，TOP SKU）", fontsize=12, pad=10)
    plt.colorbar(im, ax=ax, label="提升度")
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def draw_four_table_compare(filepath):
    """四表对比：规则数、平均提升度、平均支持度"""
    rows = []
    for name, sheet in [("⑤订单-产品线", "⑤_订单产品线"), ("⑥订单-SPU", "⑥_订单SPU"), ("⑦客户-SPU", "⑦_客户SPU"), ("⑧客户-SKU", "⑧_客户SKU")]:
        df = pd.read_excel(BASKET_EXCEL, sheet_name=sheet)
        rows.append({
            "表": name,
            "规则数": len(df),
            "平均提升度": df["提升度"].mean(),
            "平均支持度": df["支持度"].mean(),
        })
    tb = pd.DataFrame(rows)
    x = np.arange(4)
    width = 0.35
    fig, ax1 = plt.subplots(figsize=(9, 5))
    ax1.bar(x - width/2, tb["规则数"], width, label="规则数", color=MCK["dark_blue"])
    ax1.set_ylabel("规则数", fontsize=11)
    ax1.set_xticks(x)
    ax1.set_xticklabels(tb["表"])
    ax2 = ax1.twinx()
    ax2.bar(x + width/2, tb["平均提升度"], width, label="平均提升度", color=MCK["positive"], alpha=0.8)
    ax2.set_ylabel("平均提升度", fontsize=11)
    ax2.axhline(1, color="gray", linestyle="--")
    ax1.legend(loc="upper left")
    ax2.legend(loc="upper right")
    ax1.set_title("购物篮四表对比：规则数与平均提升度", fontsize=12, pad=10)
    fig.patch.set_facecolor("white")
    ax1.set_facecolor("white")
    plt.tight_layout()
    plt.savefig(filepath, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()


def add_slide_section(prs, title):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(1.2))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(5, 28, 44)
    tb.text_frame.paragraphs[0].alignment = 1


def add_slide_chart_simple(prs, title, img_path, footnote=""):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.55))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0, 0, 0)
    slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.0), width=Inches(12), height=Inches(5.5))
    if footnote:
        ft = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12), Inches(0.35))
        ft.text_frame.paragraphs[0].add_run().text = footnote
        ft.text_frame.paragraphs[0].font.size = Pt(9)
        ft.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)


def add_slide_conclusion(prs):
    """结论要点页：购物篮洞察与行动建议"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
    title.text_frame.paragraphs[0].add_run().text = "专题五 结论与行动建议"
    title.text_frame.paragraphs[0].font.size = Pt(20)
    title.text_frame.paragraphs[0].font.bold = True
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(5.5))
    tf = body.text_frame
    bullets = [
        "正向关联（提升度>1）：智能母婴电器↔吸奶器、家居家纺↔喂养电器、护理↔护理电器等，做场景化套装与推荐；订单-SPU 做加购/套装，客户-SPU 做复购与跨品拉新。",
        "互斥（提升度<1）：吸奶器与家居家纺/内衣服饰/喂养电器等同单较少，反映不同购买阶段或不同购物篮；不宜同单强绑，可分阶段、分人群触达。",
        "订单维度(⑤⑥)侧重同单搭配与满减；客户维度(⑦⑧)侧重复购、跨单推荐与 SKU 级规格/颜色推荐。",
        "被忽视组合：高提升度、中低支持度的组合可小流量测试（如冻奶壶+暖奶瓶、吸奶器+洗涤块），验证后放大。",
    ]
    for i, line in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ""
        r = p.add_run()
        r.text = "• " + line
        r.font.size = Pt(11)
        r.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(14)
    tf.word_wrap = True


def add_slide_table_scenario(prs):
    """典型组合↔典型需求场景 表格页"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.55))
    title.text_frame.paragraphs[0].add_run().text = "典型组合 ↔ 典型需求场景"
    title.text_frame.paragraphs[0].font.size = Pt(18)
    title.text_frame.paragraphs[0].font.bold = True
    # 表头 + 6 行数据
    rows_data = 7
    cols_data = 4
    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(11)
    height = Inches(0.35 * rows_data)
    table = slide.shapes.add_table(rows_data, cols_data, left, top, width, height).table
    headers = ["组合类型", "典型组合（示例）", "需求场景", "建议动作"]
    data = [
        ["订单-产品线", "智能母婴电器 ↔ 吸奶器", "新生儿家庭一站式备货", "首页/分类「新生儿必备」套装、满减"],
        ["订单-产品线", "家居家纺 ↔ 喂养电器", "居家喂养与睡眠", "「居家喂养」场景页、组合优惠"],
        ["订单-SPU", "BS03 洗涤块 ↔ M9 吸奶器", "吸奶器+清洗消毒", "加购推荐、套装、详情页推荐"],
        ["订单-SPU", "冻奶壶 ↔ MW05 暖奶瓶", "储奶与温奶", "背奶/外出场景页、组合"],
        ["客户-SPU", "先 M9 后 BS03（复购）", "耗材复购", "复购提醒、订阅/周期购"],
        ["客户-SKU", "同 SPU 不同色/规格", "款式扩充与复购", "详情页「还买了」、购物车推荐"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
            table.cell(r, c).text_frame.paragraphs[0].font.size = Pt(9)
    table.columns[0].width = Inches(1.4)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(4.0)


def main():
    if not BASKET_EXCEL.exists():
        print("请先运行 run_专题五_购物篮关联分析.py 生成", BASKET_EXCEL)
        return
    p1 = IMG_DIR / "sheet5_product_line_top_lift.png"
    p2 = IMG_DIR / "sheet5_lift_ge1_vs_lt1.png"
    p3 = IMG_DIR / "sheet5_spu_top_lift.png"
    p4 = IMG_DIR / "sheet5_four_table_compare.png"
    p5 = IMG_DIR / "sheet5_heatmap_customer_spu.png"
    p6 = IMG_DIR / "sheet5_heatmap_customer_sku.png"
    draw_product_line_top_lift(p1)
    draw_lift_ge1_vs_lt1(p2)
    draw_spu_top_lift(p3)
    draw_four_table_compare(p4)
    draw_heatmap_customer_spu(p5)
    draw_heatmap_customer_sku(p6)
    print("图表已生成:", p1, p2, p3, p4, p5, p6)
    if not OUT_PPT.exists():
        print("未找到", OUT_PPT)
        return
    prs = Presentation(str(OUT_PPT))
    remove_sheet5_slides_if_present(prs)
    add_slide_section(prs, "专题五：购物篮关联分析（⑤⑥⑦⑧）")
    add_slide_chart_simple(prs, "⑤ 订单-产品线 关联规则 TOP10（提升度≥1）", p1,
        "先继项→后继项 | 提升度>1 表示正关联 | 数据来源: 专题一 ⑤")
    add_slide_chart_simple(prs, "⑤ 规则分布：正关联(提升度≥1) vs 互斥(提升度<1)", p2,
        "提升度<1 表示先继项与后继项在订单内互斥或替代 | 数据来源: 专题一 ⑤")
    add_slide_chart_simple(prs, "⑥ 订单-SPU 关联规则 TOP10（按提升度）", p3,
        "订单维度 SPU 搭配，可用于加购/套装推荐 | 数据来源: 专题一 ⑥")
    add_slide_chart_simple(prs, "购物篮四表对比：规则数与平均提升度", p4,
        "⑤订单-产品线 ⑥订单-SPU ⑦客户-SPU ⑧客户-SKU | 数据来源: 专题一 ⑤⑥⑦⑧")
    add_slide_chart_simple(prs, "⑦ 客户-SPU 关联提升度热力图（TOP SPU）", p5,
        "行=先继项 SPU，列=后继项 SPU，色块=提升度 | 数据来源: 专题一 ⑦")
    add_slide_chart_simple(prs, "⑧ 客户-SKU 关联提升度热力图（TOP SKU）", p6,
        "行=先继项 SKU，列=后继项 SKU，色块=提升度 | 数据来源: 专题一 ⑧")
    add_slide_table_scenario(prs)
    add_slide_conclusion(prs)
    prs.save(OUT_PPT)
    print("已追加专题五 8 页（6 图 + 1 表格 + 1 结论）至:", OUT_PPT)


if __name__ == "__main__":
    main()
