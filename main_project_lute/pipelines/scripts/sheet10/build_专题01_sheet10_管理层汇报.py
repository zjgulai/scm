# -*- coding: utf-8 -*-
"""
专题一 Sheet10 用户多维洞察 - 管理层汇报 PPT（新版）

输入：
- data_example/outputs/专题一_Sheet10_用户多维洞察_计算结果.xlsx

输出：
- data_example/outputs/专题一_Sheet10_管理层汇报.pptx
"""

from __future__ import annotations

import os
from pathlib import Path

os.environ["MPLCONFIGDIR"] = str(Path(__file__).resolve().parent / ".mplcache")
os.environ["XDG_CACHE_HOME"] = str(Path(__file__).resolve().parent / ".cache")

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from matplotlib.ticker import PercentFormatter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


plt.rcParams["font.sans-serif"] = [
    "PingFang SC",
    "Heiti SC",
    "STHeiti",
    "SimHei",
    "Microsoft YaHei",
    "DejaVu Sans",
]
plt.rcParams["axes.unicode_minus"] = False


BASE_DIR = Path(__file__).resolve().parents[2]
TOPIC_DIR = BASE_DIR / "专题产物" / "专题01-sheet10"
RESULT_XLSX = TOPIC_DIR / "表格" / "专题01-sheet10_表格_用户多维洞察计算结果.xlsx"
OUT_PPT = TOPIC_DIR / "汇报" / "专题01-sheet10_汇报_管理层汇报.pptx"
IMG_DIR = TOPIC_DIR / "图表"
IMG_DIR.mkdir(parents=True, exist_ok=True)

BASE_WINDOW = "YTD2026P2"
COMPARE_WINDOW = "YTD2025P2"

MCK = {
    "dark_blue": "#051C2C",
    "medium_blue": "#1E3A5F",
    "light_blue": "#4F81BD",
    "accent_blue": "#0089FF",
    "teal": "#00A6A6",
    "green": "#27AE60",
    "orange": "#E67E22",
    "red": "#E74C3C",
    "gray_dark": "#333333",
    "gray": "#7F8C8D",
    "gray_light": "#D9E2F3",
}


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def load_data() -> dict[str, pd.DataFrame]:
    if not RESULT_XLSX.exists():
        raise FileNotFoundError(f"未找到输入文件: {RESULT_XLSX}")

    return {
        "detail": pd.read_excel(RESULT_XLSX, sheet_name="指标明细"),
        "compare": pd.read_excel(RESULT_XLSX, sheet_name="同比变动"),
        "resource": pd.read_excel(RESULT_XLSX, sheet_name="资源匹配度"),
        "promo": pd.read_excel(RESULT_XLSX, sheet_name="推广费结构"),
        "cac_retention": pd.read_excel(RESULT_XLSX, sheet_name="CAC与复购"),
    }


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_content_slide(prs: Presentation, title: str, lines: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)


def add_chart_slide(prs: Presentation, title: str, img_path: Path, note: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tr = title_box.text_frame.paragraphs[0].add_run()
    tr.text = title
    tr.font.size = Pt(20)
    tr.font.bold = True
    tr.font.color.rgb = rgb(MCK["dark_blue"])

    slide.shapes.add_picture(str(img_path), Inches(0.6), Inches(0.9), width=Inches(11.8), height=Inches(5.9))

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.75), Inches(11.8), Inches(0.4))
        nr = note_box.text_frame.paragraphs[0].add_run()
        nr.text = note
        nr.font.size = Pt(10)
        nr.font.color.rgb = rgb(MCK["gray"])


def chart_overall_kpis(detail: pd.DataFrame) -> Path:
    scope = detail[(detail["业务区域"] == "整体") & (detail["客户类型"] == "整体")]
    row_base = scope[scope["时间窗口"] == BASE_WINDOW].iloc[0]
    row_comp = scope[scope["时间窗口"] == COMPARE_WINDOW].iloc[0]

    metrics = ["ARPU值", "总推广费", "前台毛利率", "复购率"]
    labels = ["ARPU", "总推广费", "前台毛利率", "复购率"]

    base_vals = [row_base[m] for m in metrics]
    comp_vals = [row_comp[m] for m in metrics]

    fig, axes = plt.subplots(2, 2, figsize=(11.8, 6.2))
    axes = axes.flatten()
    for i, ax in enumerate(axes):
        ax.bar([COMPARE_WINDOW, BASE_WINDOW], [comp_vals[i], base_vals[i]], color=[MCK["gray_light"], MCK["accent_blue"]])
        ax.set_title(labels[i], fontsize=11)
        if labels[i] in ["前台毛利率", "复购率"]:
            ax.yaxis.set_major_formatter(PercentFormatter(1.0))
        for x, v in zip([COMPARE_WINDOW, BASE_WINDOW], [comp_vals[i], base_vals[i]]):
            if labels[i] in ["前台毛利率", "复购率"]:
                ax.text(x, v, f"{v*100:.2f}%", ha="center", va="bottom", fontsize=9)
            else:
                ax.text(x, v, f"{v:,.0f}", ha="center", va="bottom", fontsize=9)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)

    fig.suptitle("整体KPI同比对比（YTD2026P2 vs YTD2025P2）", fontsize=14, fontweight="bold")
    plt.tight_layout()
    out = IMG_DIR / "sheet10_overall_kpis.png"
    plt.savefig(out, dpi=150, bbox_inches="tight")
    plt.close()
    return out


def chart_region_compare(detail: pd.DataFrame) -> Path:
    scope = detail[(detail["时间窗口"] == BASE_WINDOW) & (detail["客户类型"] == "整体") & (detail["业务区域"].isin(["北美", "欧洲"]))]

    metrics = ["ARPU值", "人均推广费", "人均毛利额"]
    labels = ["ARPU值", "人均推广费", "人均毛利额"]

    fig, ax = plt.subplots(figsize=(11.5, 6.0))
    x = range(len(metrics))
    na_vals = [scope[scope["业务区域"] == "北美"][m].iloc[0] for m in metrics]
    eu_vals = [scope[scope["业务区域"] == "欧洲"][m].iloc[0] for m in metrics]

    width = 0.35
    ax.bar([i - width / 2 for i in x], na_vals, width=width, color=MCK["medium_blue"], label="北美")
    ax.bar([i + width / 2 for i in x], eu_vals, width=width, color=MCK["teal"], label="欧洲")

    ax.set_xticks(list(x))
    ax.set_xticklabels(labels)
    ax.set_title(f"区域对比（{BASE_WINDOW}，整体客户）", fontsize=14, fontweight="bold")
    ax.legend()
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    for i, v in enumerate(na_vals):
        ax.text(i - width / 2, v, f"{v:,.0f}", ha="center", va="bottom", fontsize=9)
    for i, v in enumerate(eu_vals):
        ax.text(i + width / 2, v, f"{v:,.0f}", ha="center", va="bottom", fontsize=9)

    plt.tight_layout()
    out = IMG_DIR / "sheet10_region_customer.png"
    plt.savefig(out, dpi=150, bbox_inches="tight")
    plt.close()
    return out


def chart_customer_compare(detail: pd.DataFrame) -> Path:
    scope = detail[(detail["时间窗口"] == BASE_WINDOW) & (detail["业务区域"] == "整体") & (detail["客户类型"].isin(["新客", "老客"]))]

    metrics = ["ARPU值", "人均毛利额", "人均订单数"]
    labels = ["ARPU值", "人均毛利额", "人均订单数"]

    fig, ax = plt.subplots(figsize=(11.5, 6.0))
    x = range(len(metrics))
    new_vals = [scope[scope["客户类型"] == "新客"][m].iloc[0] for m in metrics]
    old_vals = [scope[scope["客户类型"] == "老客"][m].iloc[0] for m in metrics]

    width = 0.35
    ax.bar([i - width / 2 for i in x], new_vals, width=width, color=MCK["orange"], label="新客")
    ax.bar([i + width / 2 for i in x], old_vals, width=width, color=MCK["green"], label="老客")

    ax.set_xticks(list(x))
    ax.set_xticklabels(labels)
    ax.set_title(f"新老客对比（{BASE_WINDOW}，整体区域）", fontsize=14, fontweight="bold")
    ax.legend()
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    for i, v in enumerate(new_vals):
        ax.text(i - width / 2, v, f"{v:,.2f}" if i == 2 else f"{v:,.0f}", ha="center", va="bottom", fontsize=9)
    for i, v in enumerate(old_vals):
        ax.text(i + width / 2, v, f"{v:,.2f}" if i == 2 else f"{v:,.0f}", ha="center", va="bottom", fontsize=9)

    plt.tight_layout()
    out = IMG_DIR / "sheet10_customer_compare.png"
    plt.savefig(out, dpi=150, bbox_inches="tight")
    plt.close()
    return out


def chart_resource_quadrant(resource: pd.DataFrame) -> Path:
    scope = resource[resource["时间窗口"] == BASE_WINDOW].copy()

    fig, ax = plt.subplots(figsize=(11.5, 6.0))
    x = scope["投入占比"]
    y = scope["产出占比"]

    ax.scatter(x, y, s=180, color=MCK["accent_blue"])
    for _, row in scope.iterrows():
        label = f"{row['业务区域']}-{row['客户类型']}"
        ax.annotate(label, (row["投入占比"], row["产出占比"]), textcoords="offset points", xytext=(6, 6), fontsize=9)

    ax.axvline(0.25, color=MCK["gray"], linestyle="--", linewidth=1)
    ax.axhline(0.25, color=MCK["gray"], linestyle="--", linewidth=1)
    ax.set_xlim(0, max(0.45, x.max() * 1.15))
    ax.set_ylim(0, max(0.45, y.max() * 1.15))
    ax.xaxis.set_major_formatter(PercentFormatter(1.0))
    ax.yaxis.set_major_formatter(PercentFormatter(1.0))
    ax.set_xlabel("投入占比")
    ax.set_ylabel("产出占比")
    ax.set_title(f"资源匹配四象限（{BASE_WINDOW}）", fontsize=14, fontweight="bold")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    plt.tight_layout()
    out = IMG_DIR / "sheet10_resource_quadrant.png"
    plt.savefig(out, dpi=150, bbox_inches="tight")
    plt.close()
    return out


def chart_cac_discount(cac_retention: pd.DataFrame, promo: pd.DataFrame) -> Path:
    cac = cac_retention[cac_retention["客户类型"] == "CAC"].copy()
    cac = cac.pivot(index="业务区域", columns="时间窗口", values="指标值").reset_index()

    discount = promo[(promo["客户类型"] == "整体") & (promo["业务区域"].isin(["北美", "欧洲"]))][["业务区域", "时间窗口", "折扣力度"]]
    discount = discount.pivot(index="业务区域", columns="时间窗口", values="折扣力度").reset_index()

    fig, axes = plt.subplots(1, 2, figsize=(12, 5.5))

    # CAC
    ax = axes[0]
    regions = cac["业务区域"].tolist()
    x = range(len(regions))
    width = 0.35
    comp_vals = [cac[cac["业务区域"] == r][COMPARE_WINDOW].iloc[0] for r in regions]
    base_vals = [cac[cac["业务区域"] == r][BASE_WINDOW].iloc[0] for r in regions]
    ax.bar([i - width / 2 for i in x], comp_vals, width=width, color=MCK["gray_light"], label=COMPARE_WINDOW)
    ax.bar([i + width / 2 for i in x], base_vals, width=width, color=MCK["orange"], label=BASE_WINDOW)
    ax.set_xticks(list(x))
    ax.set_xticklabels(regions)
    ax.set_title("CAC 同比")
    ax.legend(fontsize=8)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    # 折扣力度
    ax2 = axes[1]
    d_regions = discount["业务区域"].tolist()
    x2 = range(len(d_regions))
    comp_d = [discount[discount["业务区域"] == r][COMPARE_WINDOW].iloc[0] for r in d_regions]
    base_d = [discount[discount["业务区域"] == r][BASE_WINDOW].iloc[0] for r in d_regions]
    ax2.bar([i - width / 2 for i in x2], comp_d, width=width, color=MCK["gray_light"], label=COMPARE_WINDOW)
    ax2.bar([i + width / 2 for i in x2], base_d, width=width, color=MCK["teal"], label=BASE_WINDOW)
    ax2.set_xticks(list(x2))
    ax2.set_xticklabels(d_regions)
    ax2.set_title("折扣力度同比（|折扣额|/销售额）")
    ax2.yaxis.set_major_formatter(PercentFormatter(1.0))
    ax2.legend(fontsize=8)
    ax2.spines["top"].set_visible(False)
    ax2.spines["right"].set_visible(False)

    fig.suptitle("CAC 与折扣力度双护栏", fontsize=14, fontweight="bold")
    plt.tight_layout()
    out = IMG_DIR / "sheet10_cac_discount.png"
    plt.savefig(out, dpi=150, bbox_inches="tight")
    plt.close()
    return out


def chart_promo_mix(promo: pd.DataFrame) -> Path:
    scope = promo[(promo["时间窗口"] == BASE_WINDOW) & (promo["客户类型"] == "整体") & (promo["业务区域"].isin(["北美", "欧洲"]))]

    regions = scope["业务区域"].tolist()
    ad_share = scope["广告费占比"].tolist()
    gtm_share = scope["GTM推广费占比"].tolist()

    fig, ax = plt.subplots(figsize=(11.5, 5.8))
    ax.bar(regions, ad_share, color=MCK["medium_blue"], label="广告费占比")
    ax.bar(regions, gtm_share, bottom=ad_share, color=MCK["orange"], label="GTM推广费占比")
    ax.yaxis.set_major_formatter(PercentFormatter(1.0))
    ax.set_title(f"推广费结构（{BASE_WINDOW}，整体客户）", fontsize=14, fontweight="bold")
    ax.legend()
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    for i, r in enumerate(regions):
        ax.text(i, ad_share[i] / 2, f"{ad_share[i]*100:.1f}%", ha="center", va="center", color="white", fontsize=9)
        ax.text(i, ad_share[i] + gtm_share[i] / 2, f"{gtm_share[i]*100:.1f}%", ha="center", va="center", color="white", fontsize=9)

    plt.tight_layout()
    out = IMG_DIR / "sheet10_promo_mix.png"
    plt.savefig(out, dpi=150, bbox_inches="tight")
    plt.close()
    return out


def build_summary_lines(data: dict[str, pd.DataFrame]) -> list[str]:
    detail = data["detail"]
    resource = data["resource"]

    total_base = detail[(detail["时间窗口"] == BASE_WINDOW) & (detail["业务区域"] == "整体") & (detail["客户类型"] == "整体")].iloc[0]
    total_comp = detail[(detail["时间窗口"] == COMPARE_WINDOW) & (detail["业务区域"] == "整体") & (detail["客户类型"] == "整体")].iloc[0]

    best = resource[resource["时间窗口"] == BASE_WINDOW].sort_values("每元推广费产出毛利额", ascending=False).iloc[0]
    worst = resource[resource["时间窗口"] == BASE_WINDOW].sort_values("资源匹配差").iloc[0]

    arpu_yoy = (total_base["ARPU值"] - total_comp["ARPU值"]) / total_comp["ARPU值"]
    promo_yoy = (total_base["总推广费"] - total_comp["总推广费"]) / total_comp["总推广费"]

    return [
        f"1) 整体ARPU同比：{arpu_yoy*100:.2f}%（{total_comp['ARPU值']:.2f} -> {total_base['ARPU值']:.2f}）",
        f"2) 整体总推广费同比：{promo_yoy*100:.2f}%（{total_comp['总推广费']:.0f} -> {total_base['总推广费']:.0f}）",
        f"3) 最优效率维度：{best['业务区域']}-{best['客户类型']}，每元推广费产出毛利额 {best['每元推广费产出毛利额']:.2f}",
        f"4) 最需优化维度：{worst['业务区域']}-{worst['客户类型']}，资源匹配差 {worst['资源匹配差']*100:.2f}pp",
        "5) 建议：放大高效率维度预算，压缩负匹配维度投放，联动CAC与折扣力度双护栏治理。",
    ]


def main() -> None:
    data = load_data()

    img_overall = chart_overall_kpis(data["detail"])
    img_region = chart_region_compare(data["detail"])
    img_customer = chart_customer_compare(data["detail"])
    img_quadrant = chart_resource_quadrant(data["resource"])
    img_cac_discount = chart_cac_discount(data["cac_retention"], data["promo"])
    img_mix = chart_promo_mix(data["promo"])

    prs = Presentation()

    add_title_slide(
        prs,
        "Sheet10 用户多维洞察（新版）",
        "YTD2026P2 vs YTD2025P2 | 资源匹配度与价值产出效率",
    )

    add_chart_slide(prs, "P2 整体 KPI YoY 对比", img_overall, "口径：总推广费=广告费+GTM推广费，不含折扣额")
    add_chart_slide(prs, "P3 区域对比（北美 vs 欧洲）", img_region)
    add_chart_slide(prs, "P4 新老客对比（整体区域）", img_customer)
    add_chart_slide(prs, "P5 资源匹配四象限", img_quadrant)
    add_chart_slide(prs, "P6 CAC 同比 + 折扣力度同比", img_cac_discount)
    add_chart_slide(prs, "P7 推广费结构（广告 vs GTM）", img_mix)

    add_content_slide(prs, "P8 关键发现与行动建议", build_summary_lines(data))

    prs.save(OUT_PPT)
    print(f"PPT 已生成: {OUT_PPT}")


if __name__ == "__main__":
    main()
