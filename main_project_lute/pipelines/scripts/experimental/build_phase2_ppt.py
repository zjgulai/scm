# MOVED TO: scripts/phase2/build_phase2_ppt.py
# -*- coding: utf-8 -*-
"""
Phase2 MVP 汇报PPT生成

基于专题② + 交叉线3的分析结果，生成管理层汇报PPT
依赖: pandas, matplotlib, python-pptx
"""

import os
from pathlib import Path
import pandas as pd
import numpy as np

# 字体配置
_BASE = Path(__file__).resolve().parent
os.environ["MPLCONFIGDIR"] = str(_BASE / ".mplcache")

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
plt.rcParams["font.sans-serif"] = ["PingFang SC", "Heiti SC", "STHeiti", "SimHei", "Microsoft YaHei"]
plt.rcParams["axes.unicode_minus"] = False

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 路径配置
PROJECT_ROOT = Path(__file__).resolve().parents[2]
PHASE2_OUT = PROJECT_ROOT / "phase2_outputs" / "topic2"
CROSSLINE_OUT = PROJECT_ROOT / "phase2_outputs" / "crossline3_voc"
OUT_DIR = PROJECT_ROOT / "phase2_outputs"
OUT_PPT = OUT_DIR / "Phase2_MVP_汇报.pptx"

# 麦肯锡色板
MCK = {
    "dark_blue": "#051C2C",
    "medium_blue": "#1E3A5F",
    "light_blue": "#3A5F8A",
    "pale_blue": "#5A7F9F",
    "positive": "#27AE60",
    "negative": "#E74C3C",
    "orange": "#F39C12",
    "light_gray": "#F5F5F5",
}


def hex_to_rgb(hex_color):
    """HEX转RGB"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle

    # 样式
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(MCK["dark_blue"])


def add_content_slide(prs, title, content_lines):
    """添加内容页"""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for line in content_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.space_after = Pt(8)


def add_chart_slide(prs, title, chart_path, description=""):
    """添加图表页"""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_shape.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(MCK["dark_blue"])

    # 图表
    if chart_path and Path(chart_path).exists():
        slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1), width=Inches(9))

    # 描述
    if description:
        desc_shape = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
        tf = desc_shape.text_frame
        p = tf.add_paragraph()
        p.text = description
        p.font.size = Pt(12)
        p.font.color.rgb = hex_to_rgb(MCK["pale_blue"])


def generate_cost_attribution_chart():
    """生成成本归因图表"""
    csv_path = PHASE2_OUT / "topic2_cost_attribution.csv"
    if not csv_path.exists():
        return None

    df = pd.read_csv(csv_path)

    fig, ax = plt.subplots(figsize=(10, 6))

    # 按平台分组
    platforms = df["channel_id"].unique()
    x = np.arange(len(platforms))
    width = 0.35

    # 数据准备
    front_total = df.groupby("channel_id")["cost_front_total"].sum().values
    back_total = df.groupby("channel_id")["cost_back_total"].sum().values

    bars1 = ax.bar(x - width/2, front_total, width, label="前台成本", color=MCK["negative"])
    bars2 = ax.bar(x + width/2, back_total, width, label="后台成本", color=MCK["light_blue"])

    ax.set_xlabel("平台")
    ax.set_ylabel("成本金额 ($)")
    ax.set_title("平台成本结构：前台 vs 后台")
    ax.set_xticks(x)
    ax.set_xticklabels(platforms)
    ax.legend()

    plt.tight_layout()
    chart_path = OUT_DIR / "chart_cost_attribution.png"
    plt.savefig(chart_path, dpi=150)
    plt.close()
    return chart_path


def generate_refund_analysis_chart():
    """生成退款原因分析图表"""
    csv_path = PHASE2_OUT / "topic2_refund_attribution.csv"
    if not csv_path.exists():
        return None

    df = pd.read_csv(csv_path)

    fig, ax = plt.subplots(figsize=(10, 6))

    # 按退款原因聚合
    reason_agg = df.groupby("return_reason_code")["refund_amt"].sum().sort_values(ascending=True)

    colors = [MCK["negative"] if v > reason_agg.mean() else MCK["orange"] for v in reason_agg.values]

    ax.barh(reason_agg.index, reason_agg.values, color=colors)
    ax.set_xlabel("退款金额 ($)")
    ax.set_title("退款原因分析")

    plt.tight_layout()
    chart_path = OUT_DIR / "chart_refund_analysis.png"
    plt.savefig(chart_path, dpi=150)
    plt.close()
    return chart_path


def generate_leadtime_chart():
    """生成订单耗时图表"""
    csv_path = PHASE2_OUT / "topic2_leadtime_diagnostics.csv"
    if not csv_path.exists():
        return None

    df = pd.read_csv(csv_path)

    fig, ax = plt.subplots(figsize=(10, 6))

    # 按仓库聚合
    warehouse_agg = df.groupby("dest_warehouse")["lead_time_total"].mean().sort_values()

    ax.barh(warehouse_agg.index, warehouse_agg.values, color=MCK["medium_blue"])
    ax.set_xlabel("平均履约时长 (天)")
    ax.set_title("各仓库履约时效")

    plt.tight_layout()
    chart_path = OUT_DIR / "chart_leadtime.png"
    plt.savefig(chart_path, dpi=150)
    plt.close()
    return chart_path


def load_conclusions():
    """加载VOC结论文本"""
    txt_path = CROSSLINE_OUT / "dual_perspective_conclusions.txt"
    if not txt_path.exists():
        return "暂无VOC结论"
    with open(txt_path, "r", encoding="utf-8") as f:
        return f.read()


def main():
    print("=== 生成 Phase2 MVP 汇报 PPT ===\n")

    prs = Presentation()

    # 1. 标题页
    add_title_slide(
        prs,
        "Phase2 MVP 汇报",
        "专题②：线上订单数量与质量提升 + 交叉线3：退款反哺VOC"
    )

    # 2. 执行摘要
    add_content_slide(
        prs,
        "执行摘要",
        [
            "✓ 专题② 四个子课题完成：",
            "  • 子课题①：订单成本与费率归因",
            "  • 子课题②：订单耗时与节点诊断",
            "  • 子课题③：订单价结构与毛利归因",
            "  • 子课题④：退款多维归因",
            "",
            "✓ 交叉线3 完成：",
            "  • 订单退款 → VOC双重视图",
            "  • 产出业务洞察与改进建议",
        ]
    )

    # 3. 成本归因
    add_content_slide(
        prs,
        "子课题①：订单成本与费率归因",
        [
            "核心发现：",
            "• 亚马逊 vs DTC 平台成本结构差异显著",
            "• 前台成本（促销+推广+退款）占比平台间差异大",
            "• 后台成本（生产+头程+仓配+佣金）相对稳定",
            "",
            "数据来源：topic2_cost_attribution.csv",
        ]
    )

    chart_path = generate_cost_attribution_chart()
    if chart_path:
        add_chart_slide(prs, "平台成本结构对比", chart_path, "前台成本 vs 后台成本")

    # 4. 订单耗时
    add_content_slide(
        prs,
        "子课题②：订单耗时与节点诊断",
        [
            "核心发现：",
            "• 各仓库平均履约时长存在差异",
            "• 部分仓库超时率较高，需关注",
            "• 订单节点各阶段耗时分布",
            "",
            "数据来源：topic2_leadtime_diagnostics.csv",
        ]
    )

    chart_path = generate_leadtime_chart()
    if chart_path:
        add_chart_slide(prs, "各仓库履约时效", chart_path)

    # 5. 退款分析
    add_content_slide(
        prs,
        "子课题④：退款多维归因",
        [
            "核心发现：",
            "• 质量问题、描述不符是主要退款原因",
            "• 物流配送问题需重点关注",
            "• 部分退比例较高，存在组合优化空间",
            "",
            "数据来源：topic2_refund_attribution.csv",
        ]
    )

    chart_path = generate_refund_analysis_chart()
    if chart_path:
        add_chart_slide(prs, "退款原因分布", chart_path)

    # 6. 交叉线3 VOC结论
    conclusions = load_conclusions()
    # 截取前2000字符
    conclusions_short = conclusions[:2000] if len(conclusions) > 2000 else conclusions

    add_content_slide(
        prs,
        "交叉线3：VOC双重视图结论",
        conclusions_short.split("\n")[:15]  # 最多15行
    )

    # 7. 驱动动作
    add_content_slide(
        prs,
        "后续驱动动作",
        [
            "1. 客服话术优化：针对高频退款原因更新FAQ",
            "2. 产品详情页：补充尺码指南、实物视频",
            "3. 退换货政策：针对部分退场景优化流程",
            "4. 产品改进：反馈给产品团队进行迭代",
            "5. 物流优化：选择更可靠的承运商",
            "",
            "下一步：",
            "• 接入真实数仓数据",
            "• 扩展至其他专题",
            "• 完善PPT产出模板",
        ]
    )

    # 保存
    prs.save(OUT_PPT)
    print(f"PPT已生成: {OUT_PPT}")


if __name__ == "__main__":
    main()
