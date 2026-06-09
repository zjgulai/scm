# MOVED TO: scripts/phase1/build_专题01_品线组合洞察页.py
# -*- coding: utf-8 -*-
"""
专题①：品线组合营销策略 — 两页洞察结论 PPT（总→分，归因→策略）
追加到 专题一_归因瀑布图.pptx 末尾
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE_DIR = Path(__file__).resolve().parents[2]
OUT_PPT = BASE_DIR / "专题产物" / "专题01" / "汇报" / "专题01_汇报_归因瀑布图.pptx"


def _slide_contains_text(slide, keyword):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if keyword in (r.text or ""):
                        return True
    return False


def remove_topic1_insight_slides_if_present(prs, max_remove=5):
    """若末尾连续页为专题①洞察结论，则移除以便重新追加"""
    removed = 0
    while prs.slides and removed < max_remove:
        last = prs.slides[-1]
        if not (_slide_contains_text(last, "品线组合营销") or _slide_contains_text(last, "行动优先级")):
            break
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])
        removed += 1


def add_slide_insight_page1(prs):
    """第1页：核心结论与故事线（总 + 归因）"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    # 标题
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.55))
    t0 = title.text_frame.paragraphs[0]
    t0.add_run().text = "专题① 品线组合营销 — 核心结论与故事线"
    t0.font.size = Pt(18)
    t0.font.bold = True
    t0.font.color.rgb = RGBColor(5, 28, 44)
    # 顶层结论（总）
    block1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12), Inches(1.35))
    tf1 = block1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.add_run().text = "【总】"
    p1.font.bold = True
    p1.font.size = Pt(11)
    p1.font.color.rgb = RGBColor(0, 137, 255)
    p1.space_after = Pt(4)
    p2 = tf1.add_paragraph()
    p2.add_run().text = "品线组合营销要在「量增」与「利稳」之间取得平衡。当前订单数驱动毛利额增长（贡献约98%），但独立站毛利率持续承压、前台成本侵蚀利润，部分SPU与品线结构不合理。策略主线：稳住利润结构（控前台、清负贡献）→ 放大增长杠杆（订单数、品单价、客品数）→ 用组合与场景提升篮件数与复购。"
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(51, 51, 51)
    p2.space_after = Pt(8)
    # 五大归因
    block2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.35), Inches(12), Inches(4.2))
    tf2 = block2.text_frame
    tf2.word_wrap = True
    h = tf2.paragraphs[0]
    h.add_run().text = "【归因】五大归因"
    h.font.bold = True
    h.font.size = Pt(11)
    h.font.color.rgb = RGBColor(0, 137, 255)
    h.space_after = Pt(6)
    bullets = [
        "① 平台与区域：亚马逊贡献、独立站拖累；独立站前台占比约51%（vs 亚马逊27%），北美+欧洲关键。",
        "② 费率结构：独立站改善毛利必须先压前台（推广/促销/退款）。",
        "③ SPU与品线：负毛利箱吃利润，箱4为利润核心，喂养/护理销高利低，高毛利箱空洞。",
        "④ 订单结构：订单数主驱动；客品数略降、前台毛利率下滑，需关联推荐与控促。",
        "⑤ 购物篮：正关联做套装与推荐，互斥组合分阶段触达，提升客品数与复购。",
    ]
    for line in bullets:
        px = tf2.add_paragraph()
        px.add_run().text = line
        px.font.size = Pt(10)
        px.font.color.rgb = RGBColor(51, 51, 51)
        px.space_after = Pt(3)


def add_slide_insight_page2(prs):
    """第2页：从归因到策略 — 行动优先级"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.55))
    title.text_frame.paragraphs[0].add_run().text = "从归因到策略 — 行动优先级"
    title.text_frame.paragraphs[0].font.size = Pt(18)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(5, 28, 44)
    # P0
    y = 0.95
    for label, body in [
        ("P0 独立站毛利修复", "控推广、控促销、控退款；复核北美/欧洲定价与折扣；先前台、后后台。"),
        ("P1 SPU与品线结构优化", "清退/收缩负毛利与核心负贡献SPU；巩固箱3～4、主推正贡献SPU；喂养/护理优化定价与结构；逐步填补高毛利箱。"),
        ("P2 订单与购物篮杠杆", "关联推荐、满件/套装、购物车营销；购物篮正关联做场景化陈列与套装；复购与跨品推荐。"),
        ("P3 竞对与执行闭环", "结合竞品品线组合与定价完善策略；落地组合选品、折扣率、场景AB test与复盘。"),
    ]:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12), Inches(1.25))
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.add_run().text = label
        p1.font.bold = True
        p1.font.size = Pt(11)
        p1.font.color.rgb = RGBColor(0, 137, 255)
        p1.space_after = Pt(2)
        p2 = tf.add_paragraph()
        p2.add_run().text = body
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(51, 51, 51)
        y += 1.45


def main():
    if not OUT_PPT.exists():
        print("未找到", OUT_PPT)
        return
    prs = Presentation(str(OUT_PPT))
    remove_topic1_insight_slides_if_present(prs)
    add_slide_insight_page1(prs)
    add_slide_insight_page2(prs)
    prs.save(OUT_PPT)
    print("已追加专题① 洞察结论 2 页至:", OUT_PPT)


if __name__ == "__main__":
    main()
